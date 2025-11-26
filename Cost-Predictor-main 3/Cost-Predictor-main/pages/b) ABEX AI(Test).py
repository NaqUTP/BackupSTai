# ======================= ABEX AI — PETRONAS Gradient UI (Full App, with Project Builder + Exports) =======================
# - Data • Model • Visualization • Predict • Results • Project Builder • Compare Projects
# - Excel/PPT export with charts, heatmaps, and lines for projects & comparisons

import io
import json
import zipfile
import requests
import numpy as np
import pandas as pd
import streamlit as st

# ML/Stats
from sklearn.impute import KNNImputer, SimpleImputer
from sklearn.model_selection import train_test_split
from sklearn.preprocessing import MinMaxScaler
from sklearn.ensemble import RandomForestRegressor, GradientBoostingRegressor
from sklearn.linear_model import Ridge, Lasso
from sklearn.svm import SVR
from sklearn.tree import DecisionTreeRegressor
from sklearn.pipeline import Pipeline
from sklearn.metrics import mean_squared_error, r2_score
from scipy.stats import linregress

# Viz in-app
import plotly.express as px
import plotly.graph_objects as go

# Viz for PPT
import matplotlib.pyplot as plt
import seaborn as sns

# PPT export
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Excel export helpers
from openpyxl.formatting.rule import ColorScaleRule
from openpyxl.chart import BarChart, LineChart, Reference
from openpyxl.utils import get_column_letter

# ---------------------------------------------------------------------------------------
# PAGE CONFIG
# ---------------------------------------------------------------------------------------
st.set_page_config(
    page_title="ABEX AI RT2025",
    page_icon="💠",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ---------------------------------------------------------------------------------------
# THEME TOKENS
# ---------------------------------------------------------------------------------------
PETRONAS = {
    "teal": "#00A19B",
    "teal_dark": "#008C87",
    "purple": "#6C4DD3",
    "white": "#FFFFFF",
    "black": "#0E1116",
    "border": "rgba(0,0,0,0.10)",
}

# ---------------------------------------------------------------------------------------
# GLOBAL CSS
# ---------------------------------------------------------------------------------------
st.markdown(
    f"""
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700;800&display=swap');

/* Use Inter for main document, but DO NOT override all children so icons stay correct */
html, body {{
  font-family: 'Inter', sans-serif;
}}

[data-testid="stAppViewContainer"] {{
  background: {PETRONAS["white"]};
  color: {PETRONAS["black"]};
  padding-top: 0.5rem;
}}

#MainMenu, footer {{ visibility: hidden; }}

/* ---------------- Sidebar ---------------- */
[data-testid="stSidebar"] {{
  background: linear-gradient(180deg, {PETRONAS["teal"]} 0%, {PETRONAS["teal_dark"]} 100%) !important;
  color: #fff !important;
  border-top-right-radius: 16px;
  border-bottom-right-radius: 16px;
  box-shadow: 0 6px 20px rgba(0,0,0,0.15);
}}
[data-testid="stSidebar"] * {{ color: #fff !important; }}

/* ========================= OPTIONAL: LIGHT TWEAK TO SIDEBAR TOGGLE ========================= */
/* We let Streamlit keep its Material Icon (no more keyboard_double_arrow_right text)
   but nudge position for a floating feel. Adjust as you like. */

[data-testid="collapsedControl"] {{
  position: fixed !important;
  top: 50% !important;
  left: 10px !important;
  transform: translateY(-50%) !important;
  z-index: 9999 !important;
}}

/* ---------------- Hero Header ---------------- */
.petronas-hero {{
  border-radius: 20px;
  padding: 28px 32px;
  margin: 6px 0 18px 0;
  color: #fff;
  background: linear-gradient(135deg, {PETRONAS["teal"]}, {PETRONAS["purple"]}, {PETRONAS["black"]});
  background-size: 200% 200%;
  animation: heroGradient 8s ease-in-out infinite, fadeIn .8s ease-in-out, heroPulse 5s ease-in-out infinite;
  box-shadow: 0 10px 24px rgba(0,0,0,.12);
}}
@keyframes heroGradient {{
  0% {{ background-position: 0% 50%; }}
  50% {{ background-position: 100% 50%; }}
  100% {{ background-position: 0% 50%; }}
}}
@keyframes fadeIn {{
  from {{ opacity: 0; transform: translateY(10px); }}
  to {{ opacity: 1; transform: translateY(0); }}
}}
@keyframes heroPulse {{
  0%   {{ box-shadow: 0 0 16px rgba(0,161,155,0.45); }}
  25%  {{ box-shadow: 0 0 26px rgba(108,77,211,0.55); }}
  50%  {{ box-shadow: 0 0 36px rgba(0,161,155,0.55); }}
  75%  {{ box-shadow: 0 0 26px rgba(108,77,211,0.55); }}
  100% {{ box-shadow: 0 0 16px rgba(0,161,155,0.45); }}
}}
.petronas-hero h1 {{ margin: 0 0 5px; font-weight: 800; letter-spacing: 0.3px; }}
.petronas-hero p {{ margin: 0; opacity: .9; font-weight: 500; }}

/* ---------------- Buttons ---------------- */
.stButton > button, .stDownloadButton > button, .petronas-button {{
  border-radius: 10px;
  padding: .6rem 1.1rem;
  font-weight: 600;
  color: #fff !important;
  border: none;
  background: linear-gradient(to right, {PETRONAS["teal"]}, {PETRONAS["purple"]});
  background-size: 200% auto;
  transition: background-position .85s ease, transform .2s ease, box-shadow .25s ease;
  text-decoration: none;
  display: inline-block;
}}
.stButton > button:hover, .stDownloadButton > button:hover, .petronas-button:hover {{
  background-position: right center;
  transform: translateY(-1px);
  box-shadow: 0 6px 16px rgba(0,0,0,0.18);
}}

/* ---------------- Tabs ---------------- */
.stTabs [role="tablist"] {{
  display: flex;
  gap: 8px;
  border-bottom: none;
  padding-bottom: 6px;
}}
.stTabs [role="tab"] {{
  background: #fff;
  color: {PETRONAS["black"]};
  border-radius: 8px;
  padding: 10px 18px;
  border: 1px solid {PETRONAS["border"]};
  font-weight: 600;
  transition: all .3s ease;
  position: relative;
}}
.stTabs [role="tab"]:hover {{
  background: linear-gradient(to right, {PETRONAS["teal"]}, {PETRONAS["purple"]});
  color: #fff;
}}
.stTabs [role="tab"][aria-selected="true"] {{
  background: linear-gradient(to right, {PETRONAS["teal"]}, {PETRONAS["purple"]});
  color: #fff;
  border-color: transparent;
  box-shadow: 0 4px 16px rgba(0,0,0,0.15);
}}
.stTabs [role="tab"][aria-selected="true"]::after {{
  content: "";
  position: absolute;
  left: 10%;
  bottom: -3px;
  width: 80%;
  height: 3px;
  background: linear-gradient(90deg, {PETRONAS["teal"]}, {PETRONAS["purple"]}, {PETRONAS["teal"]});
  background-size: 200% 100%;
  border-radius: 2px;
  animation: glowSlide 2.5s linear infinite;
}}
@keyframes glowSlide {{
  0% {{ background-position: 0% 50%; }}
  50% {{ background-position: 100% 50%; }}
  100% {{ background-position: 0% 50%; }}
}}
</style>
""",
    unsafe_allow_html=True,
)

# ---------------------------------------------------------------------------------------
# HERO HEADER
# ---------------------------------------------------------------------------------------
st.markdown(
    """
<div class="petronas-hero">
  <h1>ABEX AI RT2025</h1>
  <p>Data-driven ABEX prediction & portfolio assembly</p>
</div>
""",
    unsafe_allow_html=True,
)

# ---------------------------------------------------------------------------------------
# AUTH
# ---------------------------------------------------------------------------------------
if "authenticated" not in st.session_state:
    st.session_state.authenticated = False

APPROVED_EMAILS = st.secrets.get("emails", [])
correct_password = st.secrets.get("password", None)

if not st.session_state.authenticated:
    with st.form("login_form"):
        st.markdown("#### 🔐 Access Required", unsafe_allow_html=True)
        email = st.text_input("Email Address")
        password = st.text_input("Access Password", type="password")
        submitted = st.form_submit_button("Login")

        if submitted:
            if email in APPROVED_EMAILS and password == correct_password:
                st.session_state.authenticated = True
                st.success("✅ Access granted.")
                st.rerun()
            else:
                st.error("❌ Invalid email or password.")
    st.stop()

# ---------------------------------------------------------------------------------------
# SESSION STATE
# ---------------------------------------------------------------------------------------
if "datasets" not in st.session_state:
    st.session_state.datasets = {}
if "predictions" not in st.session_state:
    st.session_state.predictions = {}
if "processed_excel_files" not in st.session_state:
    st.session_state.processed_excel_files = set()
if "_last_metrics" not in st.session_state:
    st.session_state._last_metrics = None
if "projects" not in st.session_state:
    st.session_state.projects = {}
if "component_labels" not in st.session_state:
    st.session_state.component_labels = {}
if "best_model_name_per_dataset" not in st.session_state:
    st.session_state.best_model_name_per_dataset = {}

# ---------------------------------------------------------------------------------------
# HELPERS
# ---------------------------------------------------------------------------------------
def toast(msg, icon="✅"):
    try:
        st.toast(f"{icon} {msg}")
    except Exception:
        st.success(msg if icon == "✅" else msg)


def human_format(num, pos=None):
    try:
        num = float(num)
    except Exception:
        return str(num)
    if num >= 1e9:
        return f"{num/1e9:.1f}B"
    if num >= 1e6:
        return f"{num/1e6:.1f}M"
    if num >= 1e3:
        return f"{num/1e3:.1f}K"
    return f"{num:.0f}"


def format_with_commas(num):
    try:
        return f"{float(num):,.2f}"
    except Exception:
        return str(num)


def get_currency_symbol(df: pd.DataFrame):
    for col in df.columns:
        uc = col.upper()
        if "RM" in uc:
            return "RM"
        if "USD" in uc or "$" in col:
            return "USD"
        if "€" in col:
            return "€"
        if "£" in col:
            return "£"
    try:
        sample_vals = df.iloc[:20].astype(str).values.flatten().tolist()
        if any("RM" in v.upper() for v in sample_vals):
            return "RM"
        if any("€" in v for v in sample_vals):
            return "€"
        if any("£" in v for v in sample_vals):
            return "£"
        if any("$" in v for v in sample_vals):
            return "USD"
    except Exception:
        pass
    return ""


def cost_breakdown(
    base_pred: float,
    eprr: dict,
    sst_pct: float,
    owners_pct: float,
    cont_pct: float,
    esc_pct: float,
):
    owners_cost = round(base_pred * (owners_pct / 100.0), 2)
    sst_cost = round(base_pred * (sst_pct / 100.0), 2)
    contingency_cost = round((base_pred + owners_cost) * (cont_pct / 100.0), 2)
    escalation_cost = round((base_pred + owners_cost) * (esc_pct / 100.0), 2)
    eprr_costs = {
        k: round(base_pred * (v / 100.0), 2) for k, v in (eprr or {}).items()
    }
    grand_total = round(
        base_pred + owners_cost + contingency_cost + escalation_cost, 2
    )
    return owners_cost, sst_cost, contingency_cost, escalation_cost, eprr_costs, grand_total


def project_components_df(proj):
    """
    Build a tidy DataFrame for one project: one row per component with key cost pieces.
    """
    comps = proj.get("components", [])
    rows = []
    for c in comps:
        rows.append(
            {
                "Component": c["component_type"],
                "Dataset": c["dataset"],
                "Base ABEX": float(c["prediction"]),
                "Owner's Cost": float(c["breakdown"]["owners_cost"]),
                "Contingency": float(c["breakdown"]["contingency_cost"]),
                "Escalation": float(c["breakdown"]["escalation_cost"]),
                "SST": float(c["breakdown"]["sst_cost"]),
                "Grand Total": float(c["breakdown"]["grand_total"]),
            }
        )
    return pd.DataFrame(rows)


def create_project_excel_report_abex(project_name, proj, currency=""):
    """
    Excel report for a single project:
    - Summary sheet with heatmap + bar chart + line chart
    - Component detail sheet
    """
    output = io.BytesIO()
    comps_df = project_components_df(proj)

    if comps_df.empty:
        with pd.ExcelWriter(output, engine="openpyxl") as writer:
            pd.DataFrame(
                {"Info": [f"No components for project {project_name}"]}
            ).to_excel(writer, sheet_name="Summary", index=False)
        output.seek(0)
        return output

    total_abex = comps_df["Base ABEX"].sum()
    total_grand = comps_df["Grand Total"].sum()

    summary_df = comps_df.copy()
    summary_df.loc[len(summary_df)] = {
        "Component": "TOTAL",
        "Dataset": "",
        "Base ABEX": total_abex,
        "Owner's Cost": comps_df["Owner's Cost"].sum(),
        "Contingency": comps_df["Contingency"].sum(),
        "Escalation": comps_df["Escalation"].sum(),
        "SST": comps_df["SST"].sum(),
        "Grand Total": total_grand,
    }

    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        summary_df.to_excel(writer, sheet_name="Summary", index=False)
        wb = writer.book
        ws = writer.sheets["Summary"]

        max_row = ws.max_row
        max_col = ws.max_column

        # Heatmap-style conditional formatting for numeric columns (exclude TOTAL row)
        for col_idx in range(3, max_col + 1):  # from "Base ABEX" onwards
            col_letter = get_column_letter(col_idx)
            ws.conditional_formatting.add(
                f"{col_letter}2:{col_letter}{max_row-1}",
                ColorScaleRule(
                    start_type="percentile",
                    start_value=10,
                    start_color="FFE0F7FA",
                    mid_type="percentile",
                    mid_value=50,
                    mid_color="FF80DEEA",
                    end_type="percentile",
                    end_value=90,
                    end_color="FF00838F",
                ),
            )

        # Bar chart: Grand Total by component
        chart = BarChart()
        chart.title = "Grand Total by Component"
        data = Reference(ws, min_col=8, max_col=8, min_row=1, max_row=max_row - 1)
        cats = Reference(ws, min_col=1, min_row=2, max_row=max_row - 1)
        chart.add_data(data, titles_from_data=True)
        chart.set_categories(cats)
        chart.y_axis.title = f"Cost ({currency})".strip()
        chart.x_axis.title = "Component"
        chart.height = 10
        chart.width = 18
        ws.add_chart(chart, "J2")

        # Line chart: Base ABEX trend
        line = LineChart()
        line.title = "Base ABEX Trend"
        data_abex = Reference(
            ws, min_col=3, max_col=3, min_row=1, max_row=max_row - 1
        )
        line.add_data(data_abex, titles_from_data=True)
        line.set_categories(cats)
        line.y_axis.title = f"Base ABEX ({currency})".strip()
        line.height = 10
        line.width = 18
        ws.add_chart(line, "J20")

        # Component details
        comps_df.to_excel(writer, sheet_name="Components Detail", index=False)

    output.seek(0)
    return output


def create_project_pptx_report_abex(project_name, proj, currency=""):
    """
    PowerPoint report:
    - Title slide
    - Executive summary
    - Grand Total bar chart
    - Cost composition stacked bar chart
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = f"ABEX Project Report\n{project_name}"
    p = title.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.font.size = Pt(32)
    p.font.bold = True    # noqa: F841
    p.font.color.rgb = RGBColor(0, 161, 155)

    # Summary slide
    comps_df = project_components_df(proj)
    comps = proj.get("components", [])
    total_abex = comps_df["Base ABEX"].sum() if not comps_df.empty else 0.0
    total_grand = comps_df["Grand Total"].sum() if not comps_df.empty else 0.0

    slide = prs.slides.add_slide(prs.slides[0].slide_layout)
    title = slide.shapes.title
    title.text = "Executive Summary"

    txBox = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.5), Inches(8.6), Inches(4.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    lines = [
        f"Project: {project_name}",
        "",
        f"Total Components: {len(comps)}",
        f"Total Base ABEX: {currency} {total_abex:,.2f}",
        f"Total Grand Total: {currency} {total_grand:,.2f}",
        "",
        "Components:",
    ]
    for c in comps:
        lines.append(
            f" • {c['component_type']}: {currency} {c['breakdown']['grand_total']:,.2f}"
        )

    tf.text = "\n".join(lines)
    for para in tf.paragraphs:
        para.font.size = Pt(16)

    # Charts slides
    if not comps_df.empty:
        # Grand Total by Component
        fig, ax = plt.subplots(figsize=(7, 4))
        ax.bar(comps_df["Component"], comps_df["Grand Total"])
        ax.set_title("Grand Total by Component")
        ax.set_ylabel(f"Cost ({currency})".strip())
        ax.tick_params(axis="x", rotation=25)
        ax.grid(axis="y", linestyle="--", alpha=0.4)
        fig.tight_layout()

        img_stream = io.BytesIO()
        fig.savefig(img_stream, format="png", dpi=200, bbox_inches="tight")
        plt.close(fig)
        img_stream.seek(0)

        slide = prs.slides.add_slide(prs.slides[0].slide_layout)
        title = slide.shapes.title
        title.text = "Grand Total by Component"
        slide.shapes.add_picture(
            img_stream, Inches(0.7), Inches(1.5), width=Inches(8.6)
        )

        # Stacked cost composition
        fig2, ax2 = plt.subplots(figsize=(7, 4))
        labels = comps_df["Component"]
        base = comps_df["Base ABEX"]
        owners = comps_df["Owner's Cost"]
        cont = comps_df["Contingency"]
        esc = comps_df["Escalation"]
        sst = comps_df["SST"]

        bottom = np.zeros(len(labels))
        for vals, lab in [
            (base, "Base ABEX"),
            (owners, "Owner"),
            (cont, "Contingency"),
            (esc, "Escalation"),
            (sst, "SST"),
        ]:
            ax2.bar(labels, vals, bottom=bottom, label=lab)
            bottom += vals

        ax2.set_title("Cost Composition by Component")
        ax2.set_ylabel(f"Cost ({currency})".strip())
        ax2.tick_params(axis="x", rotation=25)
        ax2.grid(axis="y", linestyle="--", alpha=0.4)
        ax2.legend(fontsize=8, ncol=3)
        fig2.tight_layout()

        img_stream2 = io.BytesIO()
        fig2.savefig(img_stream2, format="png", dpi=200, bbox_inches="tight")
        plt.close(fig2)
        img_stream2.seek(0)

        slide2 = prs.slides.add_slide(prs.slides[0].slide_layout)
        title2 = slide2.shapes.title
        title2.text = "Cost Composition by Component"
        slide2.shapes.add_picture(
            img_stream2, Inches(0.7), Inches(1.5), width=Inches(8.6)
        )

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


def create_comparison_excel_report_abex(projects_dict, currency=""):
    """
    Excel report comparing multiple projects.
    - Projects Summary with heatmap + chart
    - One sheet per project with component detail
    """
    output = io.BytesIO()

    summary_rows = []
    for name, proj in projects_dict.items():
        dfc = project_components_df(proj)
        capex = dfc["Base ABEX"].sum() if not dfc.empty else 0.0
        owners = dfc["Owner's Cost"].sum() if not dfc.empty else 0.0
        cont = dfc["Contingency"].sum() if not dfc.empty else 0.0
        esc = dfc["Escalation"].sum() if not dfc.empty else 0.0
        sst = dfc["SST"].sum() if not dfc.empty else 0.0
        grand = dfc["Grand Total"].sum() if not dfc.empty else 0.0
        summary_rows.append(
            {
                "Project": name,
                "Components": len(proj.get("components", [])),
                "ABEX Sum": capex,
                "Owner": owners,
                "Contingency": cont,
                "Escalation": esc,
                "SST": sst,
                "Grand Total": grand,
            }
        )

    summary_df = pd.DataFrame(summary_rows)

    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        summary_df.to_excel(
            writer, sheet_name="Projects Summary", index=False
        )
        wb = writer.book
        ws = writer.sheets["Projects Summary"]

        max_row = ws.max_row
        max_col = ws.max_column

        # Heatmap-style conditional formatting
        for col_idx in range(3, max_col + 1):
            col_letter = get_column_letter(col_idx)
            ws.conditional_formatting.add(
                f"{col_letter}2:{col_letter}{max_row}",
                ColorScaleRule(
                    start_type="percentile",
                    start_value=10,
                    start_color="FFE3F2FD",
                    mid_type="percentile",
                    mid_value=50,
                    mid_color="FF90CAF9",
                    end_type="percentile",
                    end_value=90,
                    end_color="FF1565C0",
                ),
            )

        # Bar chart of Grand Total
        chart = BarChart()
        chart.title = "Grand Total by Project"
        data = Reference(ws, min_col=8, max_col=8, min_row=1, max_row=max_row)
        cats = Reference(ws, min_col=1, min_row=2, max_row=max_row)
        chart.add_data(data, titles_from_data=True)
        chart.set_categories(cats)
        chart.y_axis.title = f"Cost ({currency})".strip()
        chart.x_axis.title = "Project"
        chart.height = 10
        chart.width = 18
        ws.add_chart(chart, "J2")

        # One sheet per project
        for name, proj in projects_dict.items():
            dfc = project_components_df(proj)
            if dfc.empty:
                continue
            sheet_name = name[:31]
            dfc.to_excel(writer, sheet_name=sheet_name, index=False)

    output.seek(0)
    return output


def create_comparison_pptx_report_abex(projects_dict, currency=""):
    """
    PowerPoint comparison:
    - Title slide
    - Grand Total bar chart
    - Stacked composition chart
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "ABEX Project Comparison"
    p = title.text_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 161, 155)

    # Build summary DF
    rows = []
    for name, proj in projects_dict.items():
        dfc = project_components_df(proj)
        capex = dfc["Base ABEX"].sum() if not dfc.empty else 0.0
        owners = dfc["Owner's Cost"].sum() if not dfc.empty else 0.0
        cont = dfc["Contingency"].sum() if not dfc.empty else 0.0
        esc = dfc["Escalation"].sum() if not dfc.empty else 0.0
        sst = dfc["SST"].sum() if not dfc.empty else 0.0
        grand = dfc["Grand Total"].sum() if not dfc.empty else 0.0
        rows.append(
            {
                "Project": name,
                "ABEX Sum": capex,
                "Owner": owners,
                "Contingency": cont,
                "Escalation": esc,
                "SST": sst,
                "Grand Total": grand,
            }
        )
    df_proj = pd.DataFrame(rows)

    if not df_proj.empty:
        # Grand Total by project
        fig, ax = plt.subplots(figsize=(7, 4))
        ax.bar(df_proj["Project"], df_proj["Grand Total"])
        ax.set_title("Grand Total by Project")
        ax.set_ylabel(f"Cost ({currency})".strip())
        ax.tick_params(axis="x", rotation=25)
        ax.grid(axis="y", linestyle="--", alpha=0.4)
        fig.tight_layout()

        img_stream = io.BytesIO()
        fig.savefig(img_stream, format="png", dpi=200, bbox_inches="tight")
        plt.close(fig)
        img_stream.seek(0)

        slide = prs.slides.add_slide(prs.slides[0].slide_layout)
        title = slide.shapes.title
        title.text = "Grand Total by Project"
        slide.shapes.add_picture(
            img_stream, Inches(0.7), Inches(1.5), width=Inches(8.6)
        )

        # Stacked composition by project
        fig2, ax2 = plt.subplots(figsize=(7, 4))
        labels = df_proj["Project"]
        base = df_proj["ABEX Sum"]
        owners = df_proj["Owner"]
        cont = df_proj["Contingency"]
        esc = df_proj["Escalation"]
        sst = df_proj["SST"]

        bottom = np.zeros(len(labels))
        for vals, lab in [
            (base, "Base ABEX"),
            (owners, "Owner"),
            (cont, "Contingency"),
            (esc, "Escalation"),
            (sst, "SST"),
        ]:
            ax2.bar(labels, vals, bottom=bottom, label=lab)
            bottom += vals

        ax2.set_title("Cost Composition by Project")
        ax2.set_ylabel(f"Cost ({currency})".strip())
        ax2.tick_params(axis="x", rotation=25)
        ax2.grid(axis="y", linestyle="--", alpha=0.4)
        ax2.legend(fontsize=8, ncol=3)
        fig2.tight_layout()

        img_stream2 = io.BytesIO()
        fig2.savefig(img_stream2, format="png", dpi=200, bbox_inches="tight")
        plt.close(fig2)
        img_stream2.seek(0)

        slide2 = prs.slides.add_slide(prs.slides[0].slide_layout)
        title2 = slide2.shapes.title
        title2.text = "Cost Composition by Project"
        slide2.shapes.add_picture(
            img_stream2, Inches(0.7), Inches(1.5), width=Inches(8.6)
        )

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


# ---------------------------------------------------------------------------------------
# DATA / MODEL HELPERS
# ---------------------------------------------------------------------------------------
GITHUB_USER = "Hafizuddin-Abd-Rahman-Dev-Upstream"
REPO_NAME = "Cost-Predictor"
BRANCH = "main"
DATA_FOLDER = "pages/data_ABEX"

# ---- Automatic best model (6-model pool) -----------------------------------
MODEL_CANDIDATES = {
    "RandomForest": lambda rs=42: RandomForestRegressor(random_state=rs),
    "GradientBoosting": lambda rs=42: GradientBoostingRegressor(random_state=rs),
    "Ridge": lambda rs=42: Ridge(),
    "Lasso": lambda rs=42: Lasso(),
    "SVR": lambda rs=42: SVR(),
    "DecisionTree": lambda rs=42: DecisionTreeRegressor(random_state=rs),
}


@st.cache_data(ttl=600)
def list_csvs_from_manifest(folder_path: str):
    manifest_url = f"https://raw.githubusercontent.com/{GITHUB_USER}/{REPO_NAME}/{BRANCH}/{folder_path}/files.json"
    try:
        res = requests.get(manifest_url, timeout=10)
        res.raise_for_status()
        return res.json()
    except Exception as e:
        st.error(f"Failed to load CSV manifest: {e}")
        return []


def evaluate_model(X, y, test_size=0.2, random_state=42):
    """
    Train all candidate models on a train/test split and pick the best by R².
    Returns JSON-friendly metrics (no model objects).
    """
    Xtr, Xte, ytr, yte = train_test_split(
        X, y, test_size=test_size, random_state=random_state
    )

    rows = []
    best_name = None
    best_r2 = -np.inf
    best_rmse = None

    for name, ctor in MODEL_CANDIDATES.items():
        base_model = ctor(random_state)
        pipe = Pipeline(
            steps=[
                ("imputer", SimpleImputer(strategy="median")),
                ("scaler", MinMaxScaler()),
                ("model", base_model),
            ]
        )
        pipe.fit(Xtr, ytr)
        yhat = pipe.predict(Xte)
        rmse = float(np.sqrt(mean_squared_error(yte, yhat)))
        r2 = float(r2_score(yte, yhat))
        rows.append({"model": name, "rmse": rmse, "r2": r2})
        if r2 > best_r2:
            best_r2 = r2
            best_rmse = rmse
            best_name = name

    rows_sorted = sorted(rows, key=lambda d: d["r2"], reverse=True)
    metrics = {
        "best_model": best_name,
        "rmse": best_rmse,
        "r2": best_r2,
        "models": rows_sorted,
    }
    return metrics


def get_trained_model_for_dataset(X, y, dataset_name: str, random_state=42):
    """
    Train a pipeline with the best algorithm (per dataset) on the FULL data.
    - Uses stored best algorithm if available.
    - Otherwise runs evaluate_model() once to decide, then caches the choice.
    """
    if "best_model_name_per_dataset" not in st.session_state:
        st.session_state.best_model_name_per_dataset = {}

    best_name = st.session_state.best_model_name_per_dataset.get(dataset_name)

    if not best_name:
        metrics = evaluate_model(X, y, test_size=0.2, random_state=random_state)
        best_name = metrics.get("best_model", "RandomForest")
        st.session_state.best_model_name_per_dataset[dataset_name] = best_name
        st.session_state._last_metrics = metrics

    ctor = MODEL_CANDIDATES.get(best_name, MODEL_CANDIDATES["RandomForest"])
    base_model = ctor(random_state)
    pipe = Pipeline(
        steps=[
            ("imputer", SimpleImputer(strategy="median")),
            ("scaler", MinMaxScaler()),
            ("model", base_model),
        ]
    )
    pipe.fit(X, y)
    return pipe, best_name


def single_prediction(X, y, payload: dict, dataset_name: str = "default"):
    """
    Single prediction using the best-model pipeline for this dataset.
    """
    model_pipe, _ = get_trained_model_for_dataset(X, y, dataset_name=dataset_name)
    cols = list(X.columns)
    row = {c: np.nan for c in cols}
    for c, v in payload.items():
        if c not in row:
            continue
        try:
            row[c] = float(v) if (v is not None and str(v).strip() != "") else np.nan
        except Exception:
            row[c] = np.nan
    df_in = pd.DataFrame([row], columns=cols)
    pred = float(model_pipe.predict(df_in)[0])
    return pred


# ---------------------------------------------------------------------------------------
# NAV ROW
# ---------------------------------------------------------------------------------------
nav_c1, nav_c2, nav_c3 = st.columns([1, 1, 1])
with nav_c1:
    st.button("📤 Upload Data", key="upload_top")
with nav_c2:
    if st.button("📈 New Prediction", key="predict_top"):
        for ds in list(st.session_state.predictions.keys()):
            st.session_state.predictions[ds] = []
        st.session_state.processed_excel_files = set()
        st.session_state._last_metrics = None
        for k in list(st.session_state.keys()):
            if str(k).startswith("in_"):
                del st.session_state[k]
        toast("Ready for a new prediction.")
        st.rerun()
with nav_c3:
    st.button(
        "📥 Download All",
        disabled=True,
        help="Use the Results / Project tabs for export.",
    )

# ---------------------------------------------------------------------------------------
# TABS
# ---------------------------------------------------------------------------------------
(
    tab_data,
    tab_model,
    tab_viz,
    tab_predict,
    tab_results,
    tab_pb,
    tab_compare,
) = st.tabs(
    [
        "📁 Data",
        "⚙️ Model",
        "📈 Visualization",
        "🎯 Predict",
        "📄 Results",
        "🏗️ Project Builder",
        "🔀 Compare Projects",
    ]
)

# ===================================== DATA TAB ========================================
with tab_data:
    st.markdown(
        '<h4 style="margin:0;color:#000;">Data Sources</h4><p></p>',
        unsafe_allow_html=True,
    )
    c1, c2 = st.columns([1.2, 1])
    with c1:
        data_source = st.radio(
            "Choose data source", ["Upload CSV", "Load from Server"], horizontal=True
        )
    with c2:
        st.caption("Enterprise Storage (SharePoint)")
        data_link = (
            "https://petronas.sharepoint.com/sites/ecm_ups_coe/confidential/DFE%20Cost%20Engineering/Forms/"
            "AllItems.aspx?id=%2Fsites%2Fecm%5Fups%5Fcoe%2Fconfidential%2FDFE%20Cost%20Engineering%2F2%2ETemplate%20Tools%2F"
            "Cost%20Predictor%2FDatabase%2FABEX%20%28DDRR%29%20%2D%20RT%20Q1%202025&viewid=25092e6d%2D373d%2D41fe%2D8f6f%2D486cd8cdd5b8"
        )
        st.markdown(
            f'<a href="{data_link}" target="_blank" rel="noopener" class="petronas-button">Open Enterprise Storage</a>',
            unsafe_allow_html=True,
        )

    uploaded_files = []
    if data_source == "Upload CSV":
        uploaded_files = st.file_uploader(
            "Upload CSV files (max 200MB)", type="csv", accept_multiple_files=True
        )
    else:
        github_csvs = list_csvs_from_manifest(DATA_FOLDER)
        if github_csvs:
            selected_file = st.selectbox("Choose CSV from GitHub", github_csvs)
            if st.button("Load selected CSV"):
                raw_url = f"https://raw.githubusercontent.com/{GITHUB_USER}/{REPO_NAME}/{BRANCH}/{DATA_FOLDER}/{selected_file}"
                try:
                    df = pd.read_csv(raw_url)
                    fake = type("FakeUpload", (), {"name": selected_file})
                    uploaded_files = [fake]
                    st.session_state.datasets[selected_file] = df
                    st.session_state.predictions.setdefault(selected_file, [])
                    toast(f"Loaded from GitHub: {selected_file}")
                except Exception as e:
                    st.error(f"Error loading CSV: {e}")
        else:
            st.info("No CSV files found in GitHub folder.")

    if uploaded_files:
        for up in uploaded_files:
            if up.name not in st.session_state.datasets:
                if hasattr(up, "read"):
                    df = pd.read_csv(up)
                else:
                    df = st.session_state.datasets.get(up.name, None)
                if df is not None:
                    st.session_state.datasets[up.name] = df
                    st.session_state.predictions.setdefault(up.name, [])
        toast("Dataset(s) added.")

    st.divider()
    cA, cB, cC = st.columns([1, 1, 2])
    with cA:
        if st.button("🧹 Clear all predictions"):
            st.session_state.predictions = {
                k: [] for k in st.session_state.predictions.keys()
            }
            toast("All predictions cleared.", "🧹")
    with cB:
        if st.button("🧺 Clear processed files history"):
            st.session_state.processed_excel_files = set()
            toast("Processed files history cleared.", "🧺")
    with cC:
        if st.button("🔁 Refresh server manifest"):
            list_csvs_from_manifest.clear()
            toast("Server manifest refreshed.", "🔁")

    st.divider()

    if st.session_state.datasets:
        ds_name = st.selectbox(
            "Active dataset", list(st.session_state.datasets.keys())
        )
        df = st.session_state.datasets[ds_name]
        currency = get_currency_symbol(df)
        colA, colB, colC = st.columns([1, 1, 1])
        with colA:
            st.metric("Rows", f"{df.shape[0]:,}")
        with colB:
            st.metric("Columns", f"{df.shape[1]:,}")
        with colC:
            st.metric("Currency", f"{currency or '—'}")
        with st.expander("Preview (first 10 rows)", expanded=False):
            st.dataframe(df.head(10), use_container_width=True)
    else:
        st.info("Upload or load a dataset to proceed.")

# ===================================== MODEL TAB =======================================
with tab_model:
    if not st.session_state.datasets:
        st.info("No dataset. Go to **Data** tab to upload or load.")
    else:
        ds_name = st.selectbox(
            "Dataset for model training",
            list(st.session_state.datasets.keys()),
            key="ds_model",
        )
        df = st.session_state.datasets[ds_name]

        with st.spinner("Imputing & preparing..."):
            imputed = pd.DataFrame(
                KNNImputer(n_neighbors=5).fit_transform(df), columns=df.columns
            )
            X = imputed.iloc[:, :-1]
            y = imputed.iloc[:, -1]
            target_column = y.name

        st.markdown(
            '<h4 style="margin:0;color:#000;">Train & Evaluate</h4><p>Step 2</p>',
            unsafe_allow_html=True,
        )
        c1, c2 = st.columns([1, 3])
        with c1:
            test_size = st.slider(
                "Test size",
                0.1,
                0.5,
                0.2,
                0.05,
                help="Fraction of data used for testing",
            )
            run = st.button("Run training")
        with c2:
            st.caption(
                "Automatic best-model selection over 6 regressors (with scaling & imputation)."
            )

        if run:
            with st.spinner("Training model..."):
                metrics = evaluate_model(X, y, test_size=test_size)
            c1, c2 = st.columns(2)
            with c1:
                st.metric("RMSE (best)", f"{metrics['rmse']:,.2f}")
            with c2:
                st.metric("R² (best)", f"{metrics['r2']:.3f}")

            st.session_state._last_metrics = metrics
            st.session_state.best_model_name_per_dataset[ds_name] = metrics.get(
                "best_model"
            )

            toast("Training complete.")
            st.caption(
                f"Best model selected: **{metrics.get('best_model', 'RandomForest')}**"
            )

            # MODEL COMPARISON TABLE (6 MODELS) WITH COLOUR
            try:
                models_list = metrics.get("models", [])
                if models_list:
                    df_models = pd.DataFrame(models_list)
                    df_models["rmse"] = df_models["rmse"].astype(float)
                    df_models["r2"] = df_models["r2"].astype(float)
                    df_models = df_models.set_index("model")
                    st.markdown("##### Model comparison (6-model pool)")
                    styled = (
                        df_models.style.format(
                            {"rmse": "{:,.2f}", "r2": "{:.3f}"}
                        )
                        .background_gradient(
                            subset=["r2"], cmap="YlGn"
                        )  # greener = better R²
                        .background_gradient(
                            subset=["rmse"], cmap="OrRd_r"
                        )  # darker red = lower RMSE
                    )
                    st.dataframe(styled, use_container_width=True)
            except Exception as e:
                st.warning(f"Could not render model comparison table: {e}")

# ================================== VISUALIZATION TAB =================================
with tab_viz:
    if not st.session_state.datasets:
        st.info("No dataset. Go to **Data** tab to upload or load.")
    else:
        ds_name = st.selectbox(
            "Dataset for visualization",
            list(st.session_state.datasets.keys()),
            key="ds_viz",
        )
        df = st.session_state.datasets[ds_name]
        imputed = pd.DataFrame(
            KNNImputer(n_neighbors=5).fit_transform(df), columns=df.columns
        )
        X = imputed.iloc[:, :-1]
        y = imputed.iloc[:, -1]
        target_column = y.name

        # Correlation Matrix
        st.markdown(
            '<h4 style="margin:0;color:#000;">Correlation Matrix</h4><p>Exploration</p>',
            unsafe_allow_html=True,
        )
        corr = imputed.corr(numeric_only=True)
        fig = px.imshow(
            corr,
            text_auto=".2f",
            aspect="auto",
            color_continuous_scale="RdBu_r",
            zmin=-1,
            zmax=1,
        )
        fig.update_layout(
            margin=dict(l=0, r=0, t=10, b=0),
            paper_bgcolor=PETRONAS["white"],
            plot_bgcolor=PETRONAS["white"],
            font=dict(color=PETRONAS["black"]),
            xaxis=dict(color=PETRONAS["black"]),
            yaxis=dict(color=PETRONAS["black"]),
        )
        st.plotly_chart(fig, use_container_width=True)
        st.markdown("<div style='height:8px;'></div>", unsafe_allow_html=True)

        # Feature Importance (RandomForest for explainability only)
        st.markdown(
            '<h4 style="margin:0;color:#000;">Feature Importance</h4><p>Model</p>',
            unsafe_allow_html=True,
        )
        scaler = MinMaxScaler().fit(X)
        model = RandomForestRegressor(random_state=42).fit(
            scaler.transform(X), y
        )
        importances = model.feature_importances_
        fi = (
            pd.DataFrame(
                {"feature": X.columns, "importance": importances}
            ).sort_values("importance", ascending=True)
        )
        fig2 = go.Figure(
            go.Bar(
                x=fi["importance"],
                y=fi["feature"],
                orientation="h",
                marker_color=PETRONAS["teal"],
            )
        )
        fig2.update_layout(
            xaxis_title="Importance",
            yaxis_title="Feature",
            margin=dict(l=0, r=0, t=10, b=0),
            paper_bgcolor=PETRONAS["white"],
            plot_bgcolor=PETRONAS["white"],
            font=dict(color=PETRONAS["black"]),
            xaxis=dict(color=PETRONAS["black"]),
            yaxis=dict(color=PETRONAS["black"]),
        )
        st.plotly_chart(fig2, use_container_width=True)
        st.markdown("<div style='height:8px;'></div>", unsafe_allow_html=True)

        # Cost Curve
        st.markdown(
            '<h4 style="margin:0;color:#000;">Cost Curve</h4><p>Trend</p>',
            unsafe_allow_html=True,
        )
        feat = st.selectbox("Select feature for cost curve", X.columns)
        x_vals = imputed[feat].values
        y_vals = y.values
        mask = (~np.isnan(x_vals)) & (~np.isnan(y_vals))

        scatter_df = pd.DataFrame(
            {feat: x_vals[mask], target_column: y_vals[mask]}
        )
        fig3 = px.scatter(scatter_df, x=feat, y=target_column, opacity=0.65)
        fig3.update_traces(marker=dict(color=PETRONAS["teal"]))

        if mask.sum() >= 2 and np.unique(x_vals[mask]).size >= 2:
            xv = scatter_df[feat].to_numpy(dtype=float)
            yv = scatter_df[target_column].to_numpy(dtype=float)
            slope, intercept, r_value, p_value, std_err = linregress(xv, yv)
            x_line = np.linspace(xv.min(), xv.max(), 100)
            y_line = slope * x_line + intercept
            fig3.add_trace(
                go.Scatter(
                    x=x_line,
                    y=y_line,
                    mode="lines",
                    name=f"Fit: y={slope:.2f}x+{intercept:.2f} (R²={r_value**2:.3f})",
                    line=dict(color=PETRONAS["purple"]),
                )
            )
        else:
            st.warning("Not enough valid/variable data to compute regression.")

        fig3.update_layout(
            margin=dict(l=0, r=0, t=10, b=0),
            paper_bgcolor=PETRONAS["white"],
            plot_bgcolor=PETRONAS["white"],
            font=dict(color=PETRONAS["black"]),
            xaxis=dict(color=PETRONAS["black"]),
            yaxis=dict(color=PETRONAS["black"]),
        )
        st.plotly_chart(fig3, use_container_width=True)
        st.markdown("<div style='height:8px;'></div>", unsafe_allow_html=True)

# ===================================== PREDICT TAB =====================================
with tab_predict:
    if not st.session_state.datasets:
        st.info("No dataset. Go to **Data** tab to upload or load.")
    else:
        ds_name = st.selectbox(
            "Dataset for prediction",
            list(st.session_state.datasets.keys()),
            key="ds_pred",
        )
        df = st.session_state.datasets[ds_name]
        currency = get_currency_symbol(df)

        imputed = pd.DataFrame(
            KNNImputer(n_neighbors=5).fit_transform(df), columns=df.columns
        )
        X, y = imputed.iloc[:, :-1], imputed.iloc[:, -1]
        target_column = y.name

        st.markdown(
            '<h4 style="margin:0;color:#000;">Configuration (EPRR • Taxes • Owner • Risk)</h4><p>Step 3</p>',
            unsafe_allow_html=True,
        )
        c1, c2 = st.columns([1, 1])
        with c1:
            st.markdown("**EPRR Breakdown (%)**")
            eng = st.slider("Engineering", 0, 100, 12)
            prep = st.slider("Preparation", 0, 100, 7)
            remv = st.slider("Removal", 0, 100, 54)
            remd = st.slider("Remediation", 0, 100, 27)
        with c2:
            st.markdown("**Financial (%)**")
            sst_pct = st.slider("SST", 0, 100, 0)
            owners_pct = st.slider("Owner's Cost", 0, 100, 0)
            cont_pct = st.slider("Contingency", 0, 100, 0)
            esc_pct = st.slider("Escalation & Inflation", 0, 100, 0)

        eprr = {"Engineering": eng, "Preparation": prep, "Removal": remv, "Remediation": remd}
        eprr_total = sum(eprr.values())
        if abs(eprr_total - 100) > 1e-6 and eprr_total > 0:
            st.warning(
                f"EPRR total is {eprr_total}%. Consider normalizing to 100% for reporting consistency."
            )

        st.markdown(
            '<h4 style="margin:0;color:#000;">Predict (Single)</h4><p>Step 4</p>',
            unsafe_allow_html=True,
        )
        project_name = st.text_input(
            "Project Name",
            placeholder="e.g., Offshore Pipeline Replacement 2025",
        )
        st.caption("Provide feature values (leave blank for NaN).")

        cols_per_row = 3
        new_data = {}
        cols = list(X.columns)
        rows = (len(cols) + cols_per_row - 1) // cols_per_row
        for r in range(rows):
            row_cols = st.columns(cols_per_row)
            for i in range(cols_per_row):
                idx = r * cols_per_row + i
                if idx < len(cols):
                    col_name = cols[idx]
                    with row_cols[i]:
                        val = st.text_input(col_name, key=f"in_{col_name}")
                        new_data[col_name] = val

        if st.button("Run Prediction"):
            pred = single_prediction(X, y, new_data, dataset_name=ds_name)
            (
                owners_cost,
                sst_cost,
                contingency_cost,
                escalation_cost,
                eprr_costs,
                grand_total,
            ) = cost_breakdown(
                pred, eprr, sst_pct, owners_pct, cont_pct, esc_pct
            )

            result = {
                "Project Name": project_name,
                **{c: new_data[c] for c in cols},
                target_column: round(pred, 2),
            }
            for k, v in eprr_costs.items():
                result[f"{k} Cost"] = v
            result["SST Cost"] = sst_cost
            result["Owner's Cost"] = owners_cost
            result["Cost Contingency"] = contingency_cost
            result["Escalation & Inflation"] = escalation_cost
            result["Grand Total"] = grand_total
            st.session_state.predictions.setdefault(ds_name, []).append(result)
            toast("Prediction added to Results.")

            cA, cB, cC, cD, cE = st.columns(5)
            with cA:
                st.metric("Predicted", f"{currency} {pred:,.2f}")
            with cB:
                st.metric("Owner's", f"{currency} {owners_cost:,.2f}")
            with cC:
                st.metric("Contingency", f"{currency} {contingency_cost:,.2f}")
            with cD:
                st.metric("Escalation", f"{currency} {escalation_cost:,.2f}")
            with cE:
                st.metric("Grand Total", f"{currency} {grand_total:,.2f}")

        st.markdown(
            '<h4 style="margin:0;color:#000;">Batch (Excel)</h4>',
            unsafe_allow_html=True,
        )
        xls = st.file_uploader(
            "Upload Excel for batch prediction", type=["xlsx"]
        )
        if xls:
            file_id = f"{xls.name}_{xls.size}_{ds_name}"
            if file_id not in st.session_state.processed_excel_files:
                batch_df = pd.read_excel(xls)
                missing = [c for c in X.columns if c not in batch_df.columns]
                if missing:
                    st.error(
                        f"Missing required columns in Excel: {missing}"
                    )
                else:
                    model_pipe, best_name = get_trained_model_for_dataset(
                        X, y, dataset_name=ds_name
                    )
                    preds = model_pipe.predict(batch_df[X.columns])
                    batch_df[target_column] = preds

                    for i, row in batch_df.iterrows():
                        name = row.get("Project Name", f"Project {i+1}")
                        entry = {"Project Name": name}
                        entry.update(row[X.columns].to_dict())
                        entry[target_column] = round(float(preds[i]), 2)
                        (
                            owners_cost,
                            sst_cost,
                            contingency_cost,
                            escalation_cost,
                            eprr_costs,
                            grand_total,
                        ) = cost_breakdown(
                            float(preds[i]),
                            eprr,
                            sst_pct,
                            owners_pct,
                            cont_pct,
                            esc_pct,
                        )
                        for k, v in eprr_costs.items():
                            entry[f"{k} Cost"] = v
                        entry["SST Cost"] = sst_cost
                        entry["Owner's Cost"] = owners_cost
                        entry["Cost Contingency"] = contingency_cost
                        entry["Escalation & Inflation"] = escalation_cost
                        entry["Grand Total"] = grand_total
                        st.session_state.predictions.setdefault(
                            ds_name, []
                        ).append(entry)

                    st.session_state.processed_excel_files.add(file_id)
                    toast("Batch prediction complete.")

# ===================================== RESULTS TAB =====================================
with tab_results:
    if not st.session_state.datasets:
        st.info("No dataset. Go to **Data** tab to upload or load.")
    else:
        ds_name = st.selectbox(
            "Dataset", list(st.session_state.datasets.keys()), key="ds_results"
        )
        preds = st.session_state.predictions.get(ds_name, [])

        st.markdown(
            f'<h4 style="margin:0;color:#000;">Project Entries</h4><p>{len(preds)} saved</p>',
            unsafe_allow_html=True,
        )
        if preds:
            if st.button("🗑️ Delete all entries"):
                st.session_state.predictions[ds_name] = []
                to_remove = {
                    fid
                    for fid in st.session_state.processed_excel_files
                    if fid.endswith(ds_name)
                }
                for fid in to_remove:
                    st.session_state.processed_excel_files.remove(fid)
                toast("All entries removed.", "🗑️")
                st.rerun()

        st.markdown("<div style='height:8px;'></div>", unsafe_allow_html=True)
        st.markdown(
            '<h4 style="margin:0;color:#000;">Summary Table & Export</h4><p>Download</p>',
            unsafe_allow_html=True,
        )

        if preds:
            df_preds = pd.DataFrame(preds)
            df_disp = df_preds.copy()
            num_cols = df_disp.select_dtypes(include=[np.number]).columns
            for col in num_cols:
                df_disp[col] = df_disp[col].apply(
                    lambda x: format_with_commas(x)
                )
            st.dataframe(df_disp, use_container_width=True, height=420)

            bio_xlsx = io.BytesIO()
            df_preds.to_excel(bio_xlsx, index=False, engine="openpyxl")
            bio_xlsx.seek(0)
            metrics = st.session_state._last_metrics
            metrics_json = json.dumps(
                metrics if metrics else {"info": "No metrics"},
                indent=2,
                default=float,
            )

            zip_bio = io.BytesIO()
            with zipfile.ZipFile(
                zip_bio, "w", zipfile.ZIP_DEFLATED
            ) as zf:
                zf.writestr(
                    f"{ds_name}_predictions.xlsx", bio_xlsx.getvalue()
                )
                zf.writestr(f"{ds_name}_metrics.json", metrics_json)
            zip_bio.seek(0)

            st.download_button(
                "⬇️ Download All (ZIP)",
                data=zip_bio.getvalue(),
                file_name=f"{ds_name}_abex_all.zip",
                mime="application/zip",
            )
        else:
            st.info("No data to export yet.")

# ============================== PROJECT BUILDER TAB ====================================
with tab_pb:
    st.markdown(
        '<h4 style="margin:0;color:#000;">Project Builder</h4><p>Assemble multi-component ABEX projects</p>',
        unsafe_allow_html=True,
    )

    if not st.session_state.datasets:
        st.info("No dataset. Go to **Data** tab to upload or load.")
    else:
        colA, colB = st.columns([2, 1])
        with colA:
            new_project_name = st.text_input(
                "New Project Name",
                placeholder="e.g., ABEX Decom Campaign 2026",
                key="pb_new_project_name",
            )
        with colB:
            if new_project_name and new_project_name not in st.session_state.projects:
                if st.button("Create Project", key="pb_create_project_btn"):
                    st.session_state.projects[new_project_name] = {
                        "components": [],
                        "totals": {},
                        "currency": "",
                    }
                    toast(f"Project '{new_project_name}' created.")

        if not st.session_state.projects:
            st.info("Create a project above, then add components.")
        else:
            existing_projects = list(st.session_state.projects.keys())
            proj_sel = st.selectbox(
                "Select project to work on",
                existing_projects,
                key="pb_project_select",
            )

            ds_names = sorted(st.session_state.datasets.keys())
            dataset_for_comp = st.selectbox(
                "Dataset for this component",
                ds_names,
                key="pb_dataset_for_component",
            )

            df_comp = st.session_state.datasets[dataset_for_comp]
            currency_ds = get_currency_symbol(df_comp)

            imputed_comp = pd.DataFrame(
                KNNImputer(n_neighbors=5).fit_transform(df_comp),
                columns=df_comp.columns,
            )
            X_comp = imputed_comp.iloc[:, :-1]
            y_comp = imputed_comp.iloc[:, -1]
            target_column_comp = y_comp.name

            default_label = st.session_state.component_labels.get(
                dataset_for_comp, ""
            )
            component_type = st.text_input(
                "Component type (Asset / Scope)",
                value=(default_label or "Platform / Pipeline / Subsea / Well"),
                key=f"pb_component_type_{proj_sel}",
            )

            st.markdown("**Component Feature Inputs**")
            feat_cols = list(X_comp.columns)
            comp_inputs = {}
            cols_per_row = 2
            rows = (len(feat_cols) + cols_per_row - 1) // cols_per_row
            for r in range(rows):
                row_cols = st.columns(cols_per_row)
                for i in range(cols_per_row):
                    idx = r * cols_per_row + i
                    if idx < len(feat_cols):
                        col_name = feat_cols[idx]
                        with row_cols[i]:
                            key = f"pb_{proj_sel}_{dataset_for_comp}_feat_{col_name}"
                            comp_inputs[col_name] = st.text_input(
                                col_name, key=key
                            )

            st.markdown("---")
            st.markdown("**Cost Percentage Inputs**")
            cp1, cp2 = st.columns(2)
            with cp1:
                st.markdown("EPRR (%)")
                eng_pb = st.slider(
                    "Engineering", 0, 100, 12, key=f"pb_eng_{proj_sel}"
                )
                prep_pb = st.slider(
                    "Preparation", 0, 100, 7, key=f"pb_prep_{proj_sel}"
                )
                remv_pb = st.slider(
                    "Removal", 0, 100, 54, key=f"pb_remv_{proj_sel}"
                )
                remd_pb = st.slider(
                    "Remediation", 0, 100, 27, key=f"pb_remd_{proj_sel}"
                )
            with cp2:
                st.markdown("Financial (%)")
                sst_pb = st.slider(
                    "SST", 0, 100, 0, key=f"pb_sst_{proj_sel}"
                )
                owners_pb = st.slider(
                    "Owner's Cost", 0, 100, 0, key=f"pb_owners_{proj_sel}"
                )
                cont_pb = st.slider(
                    "Contingency", 0, 100, 0, key=f"pb_cont_{proj_sel}"
                )
                esc_pb = st.slider(
                    "Escalation & Inflation",
                    0,
                    100,
                    0,
                    key=f"pb_esc_{proj_sel}",
                )

            eprr_pb = {
                "Engineering": eng_pb,
                "Preparation": prep_pb,
                "Removal": remv_pb,
                "Remediation": remd_pb,
            }
            eprr_total_pb = sum(eprr_pb.values())
            if abs(eprr_total_pb - 100) > 1e-6 and eprr_total_pb > 0:
                st.warning(
                    f"EPRR total is {eprr_total_pb}%. Consider normalizing to 100% for reporting consistency."
                )

            if st.button(
                "➕ Predict & Add Component",
                key=f"pb_add_comp_{proj_sel}_{dataset_for_comp}",
            ):
                row_payload = {}
                for f in feat_cols:
                    v = comp_inputs.get(f, "")
                    if v is None or str(v).strip() == "":
                        row_payload[f] = np.nan
                    else:
                        try:
                            row_payload[f] = float(v)
                        except Exception:
                            row_payload[f] = np.nan
                try:
                    base_pred = single_prediction(
                        X_comp, y_comp, row_payload, dataset_name=dataset_for_comp
                    )
                    (
                        owners_cost,
                        sst_cost,
                        contingency_cost,
                        escalation_cost,
                        eprr_costs,
                        grand_total,
                    ) = cost_breakdown(
                        base_pred,
                        eprr_pb,
                        sst_pb,
                        owners_pb,
                        cont_pb,
                        esc_pb,
                    )
                    comp_entry = {
                        "component_type": component_type
                        or default_label
                        or "Component",
                        "dataset": dataset_for_comp,
                        "model_used": st.session_state.best_model_name_per_dataset.get(
                            dataset_for_comp, "BestModel"
                        ),
                        "inputs": {k: row_payload[k] for k in feat_cols},
                        "prediction": base_pred,
                        "breakdown": {
                            "eprr_costs": eprr_costs,
                            "eprr_pct": eprr_pb,
                            "sst_cost": sst_cost,
                            "owners_cost": owners_cost,
                            "contingency_cost": contingency_cost,
                            "escalation_cost": escalation_cost,
                            "grand_total": grand_total,
                            "target_col": target_column_comp,
                        },
                    }
                    st.session_state.projects[proj_sel][
                        "components"
                    ].append(comp_entry)
                    st.session_state.component_labels[
                        dataset_for_comp
                    ] = component_type or default_label
                    if not st.session_state.projects[proj_sel]["currency"]:
                        st.session_state.projects[proj_sel][
                            "currency"
                        ] = currency_ds
                    toast(f"Component added to project '{proj_sel}'.")
                except Exception as e:
                    st.error(f"Failed to predict component ABEX: {e}")

            st.markdown("---")
            st.markdown("### Current Project Overview")

            proj = st.session_state.projects[proj_sel]
            comps = proj.get("components", [])
            if not comps:
                st.info("No components yet. Add at least one above.")
            else:
                rows = []
                for c in comps:
                    rows.append(
                        {
                            "Component": c["component_type"],
                            "Dataset": c["dataset"],
                            "Model": c.get("model_used", "N/A"),
                            "Base ABEX": c["prediction"],
                            "Grand Total": c["breakdown"]["grand_total"],
                        }
                    )
                dfc = pd.DataFrame(rows)
                curr = proj.get("currency", "") or currency_ds
                st.dataframe(
                    dfc.style.format(
                        {"Base ABEX": "{:,.2f}", "Grand Total": "{:,.2f}"}
                    ),
                    use_container_width=True,
                )
                total_capex = float(sum(r["Base ABEX"] for r in rows))
                total_grand = float(sum(r["Grand Total"] for r in rows))
                proj["totals"] = {
                    "capex_sum": total_capex,
                    "grand_total": total_grand,
                }
                col_t1, col_t2 = st.columns(2)
                with col_t1:
                    st.metric("Project ABEX", f"{curr} {total_capex:,.2f}")
                with col_t2:
                    st.metric(
                        "Project Grand Total",
                        f"{curr} {total_grand:,.2f}",
                    )

                st.markdown("#### Component Cost Composition")
                comp_cost_rows = []
                for c in comps:
                    base = float(c["prediction"])
                    owners = float(c["breakdown"]["owners_cost"])
                    cont = float(c["breakdown"]["contingency_cost"])
                    esc = float(c["breakdown"]["escalation_cost"])
                    sst = float(c["breakdown"]["sst_cost"])
                    comp_cost_rows.append(
                        {
                            "Component": c["component_type"],
                            "ABEX": base,
                            "Owner": owners,
                            "Contingency": cont,
                            "Escalation": esc,
                            "SST": sst,
                        }
                    )
                df_cost = pd.DataFrame(comp_cost_rows)
                if not df_cost.empty:
                    df_melt = df_cost.melt(
                        id_vars="Component",
                        var_name="Cost Type",
                        value_name="Value",
                    )
                    fig_stack = px.bar(
                        df_melt,
                        x="Component",
                        y="Value",
                        color="Cost Type",
                        barmode="stack",
                        labels={"Value": f"Cost ({curr})"},
                        color_discrete_sequence=[
                            PETRONAS["teal"],
                            "#6C757D",
                            "#C0392B",
                            "#17A589",
                            "#9B59B6",
                        ],
                    )
                    fig_stack.update_layout(
                        margin=dict(l=0, r=0, t=10, b=0),
                        paper_bgcolor=PETRONAS["white"],
                        plot_bgcolor=PETRONAS["white"],
                        font=dict(color=PETRONAS["black"]),
                        xaxis=dict(color=PETRONAS["black"], tickangle=25),
                        yaxis=dict(color=PETRONAS["black"]),
                    )
                    st.plotly_chart(fig_stack, use_container_width=True)

                st.markdown("#### Components")
                for idx, c in enumerate(comps):
                    col1, col2, col3 = st.columns([4, 2, 1])
                    with col1:
                        st.write(
                            f"**{c['component_type']}** — *{c['dataset']}* — {c.get('model_used', 'N/A')}"
                        )
                    with col2:
                        st.write(
                            f"Grand Total: {curr} {c['breakdown']['grand_total']:,.2f}"
                        )
                    with col3:
                        if st.button(
                            "🗑️", key=f"pb_del_comp_{proj_sel}_{idx}"
                        ):
                            comps.pop(idx)
                            toast("Component removed.", "🗑️")
                            st.rerun()

                st.markdown("---")
                st.markdown("#### Export / Import Project")

                col_dl1, col_dl2, col_dl3 = st.columns(3)
                curr = proj.get("currency", "") or currency_ds

                with col_dl1:
                    excel_report = create_project_excel_report_abex(
                        proj_sel, proj, curr
                    )
                    st.download_button(
                        "⬇️ Download Project Excel",
                        data=excel_report,
                        file_name=f"{proj_sel}_ABEX_Report.xlsx",
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    )

                with col_dl2:
                    pptx_report = create_project_pptx_report_abex(
                        proj_sel, proj, curr
                    )
                    st.download_button(
                        "⬇️ Download Project PowerPoint",
                        data=pptx_report,
                        file_name=f"{proj_sel}_ABEX_Report.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    )

                with col_dl3:
                    st.download_button(
                        "⬇️ Download Project (JSON)",
                        data=json.dumps(proj, indent=2, default=float),
                        file_name=f"{proj_sel}.json",
                        mime="application/json",
                    )

                up_json = st.file_uploader(
                    "Import project JSON",
                    type=["json"],
                    key=f"pb_import_{proj_sel}",
                )
                if up_json is not None:
                    try:
                        data = json.load(up_json)
                        st.session_state.projects[proj_sel] = data
                        toast("Project imported successfully.")
                        st.rerun()
                    except Exception as e:
                        st.error(f"Failed to import project JSON: {e}")

# ============================== COMPARE PROJECTS TAB ===================================
with tab_compare:
    st.markdown(
        '<h4 style="margin:0;color:#000;">Compare Projects</h4><p>Portfolio-level ABEX view</p>',
        unsafe_allow_html=True,
    )

    proj_names = list(st.session_state.projects.keys())
    if len(proj_names) < 2:
        st.info("Create at least two projects in the Project Builder tab to compare.")
    else:
        compare_sel = st.multiselect(
            "Select projects to compare",
            proj_names,
            default=proj_names[:2],
            key="compare_projects_sel",
        )

        if len(compare_sel) < 2:
            st.warning("Select at least two projects for a meaningful comparison.")
        else:
            rows = []
            for p in compare_sel:
                proj = st.session_state.projects[p]
                comps = proj.get("components", [])
                capex = owners = cont = esc = sst = 0.0
                for c in comps:
                    capex += float(c["prediction"])
                    owners += float(c["breakdown"]["owners_cost"])
                    cont += float(c["breakdown"]["contingency_cost"])
                    esc += float(c["breakdown"]["escalation_cost"])
                    sst += float(c["breakdown"]["sst_cost"])
                grand_total = float(capex + owners + cont + esc)
                proj["totals"] = {
                    "capex_sum": capex,
                    "grand_total": grand_total,
                }
                rows.append(
                    {
                        "Project": p,
                        "Components": len(comps),
                        "ABEX Sum": capex,
                        "Owner": owners,
                        "Contingency": cont,
                        "Escalation": esc,
                        "SST": sst,
                        "Grand Total": grand_total,
                        "Currency": proj.get("currency", ""),
                    }
                )

            df_proj = pd.DataFrame(rows)
            st.dataframe(
                df_proj[
                    ["Project", "Components", "ABEX Sum", "Grand Total"]
                ].style.format(
                    {"ABEX Sum": "{:,.2f}", "Grand Total": "{:,.2f}"}
                ),
                use_container_width=True,
            )

            st.markdown("#### Grand Total by Project")
            fig_gt = px.bar(
                df_proj,
                x="Project",
                y="Grand Total",
                text="Grand Total",
                labels={"Grand Total": "Grand Total"},
                color="Project",
                color_discrete_sequence=px.colors.qualitative.Set2,
            )
            fig_gt.update_traces(
                texttemplate="%{text:,.0f}", textposition="outside"
            )
            fig_gt.update_layout(
                margin=dict(l=0, r=0, t=10, b=0),
                paper_bgcolor=PETRONAS["white"],
                plot_bgcolor=PETRONAS["white"],
                font=dict(color=PETRONAS["black"]),
                xaxis=dict(color=PETRONAS["black"], tickangle=25),
                yaxis=dict(color=PETRONAS["black"]),
            )
            st.plotly_chart(fig_gt, use_container_width=True)

            st.markdown("#### Stacked Cost Composition by Project")
            df_melt = df_proj.melt(
                id_vars=["Project"],
                value_vars=["ABEX Sum", "Owner", "Contingency", "Escalation", "SST"],
                var_name="Cost Type",
                value_name="Value",
            )
            fig_comp = px.bar(
                df_melt,
                x="Project",
                y="Value",
                color="Cost Type",
                barmode="stack",
                labels={"Value": "Cost"},
                color_discrete_sequence=[
                    PETRONAS["teal"],
                    "#6C757D",
                    "#C0392B",
                    "#17A589",
                    "#9B59B6",
                ],
            )
            fig_comp.update_layout(
                margin=dict(l=0, r=0, t=10, b=0),
                paper_bgcolor=PETRONAS["white"],
                plot_bgcolor=PETRONAS["white"],
                font=dict(color=PETRONAS["black"]),
                xaxis=dict(color=PETRONAS["black"]),
                yaxis=dict(color=PETRONAS["black"]),
            )
            st.plotly_chart(fig_comp, use_container_width=True)

            st.markdown("#### Component-Level Details")
            for p in compare_sel:
                proj = st.session_state.projects[p]
                comps = proj.get("components", [])
                if not comps:
                    continue
                with st.expander(f"Project: {p}"):
                    rows_c = []
                    for c in comps:
                        eprr_str = ", ".join(
                            f"{k}: {v:,.0f}"
                            for k, v in c["breakdown"]["eprr_costs"].items()
                            if v != 0
                        )
                        rows_c.append(
                            {
                                "Component": c["component_type"],
                                "Dataset": c["dataset"],
                                "Base ABEX": c["prediction"],
                                "Owner": c["breakdown"]["owners_cost"],
                                "Contingency": c[
                                    "breakdown"
                                ]["contingency_cost"],
                                "Escalation": c["breakdown"]["escalation_cost"],
                                "SST": c["breakdown"]["sst_cost"],
                                "Grand Total": c["breakdown"]["grand_total"],
                                "EPRR Costs": eprr_str,
                            }
                        )
                    df_compd = pd.DataFrame(rows_c)
                    st.dataframe(
                        df_compd.style.format(
                            {
                                "Base ABEX": "{:,.2f}",
                                "Owner": "{:,.2f}",
                                "Contingency": "{:,.2f}",
                                "Escalation": "{:,.2f}",
                                "SST": "{:,.2f}",
                                "Grand Total": "{:,.2f}",
                            }
                        ),
                        use_container_width=True,
                    )

            st.markdown("---")
            st.markdown("#### Download Comparison Reports")

            col_c1, col_c2 = st.columns(2)

            projects_to_export = {
                name: st.session_state.projects[name] for name in compare_sel
            }
            currency_comp = st.session_state.projects[compare_sel[0]].get(
                "currency", ""
            )

            with col_c1:
                excel_comp = create_comparison_excel_report_abex(
                    projects_to_export, currency_comp
                )
                st.download_button(
                    "⬇️ Download Comparison Excel",
                    data=excel_comp,
                    file_name="ABEX_Projects_Comparison.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                )

            with col_c2:
                pptx_comp = create_comparison_pptx_report_abex(
                    projects_to_export, currency_comp
                )
                st.download_button(
                    "⬇️ Download Comparison PowerPoint",
                    data=pptx_comp,
                    file_name="ABEX_Projects_Comparison.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
