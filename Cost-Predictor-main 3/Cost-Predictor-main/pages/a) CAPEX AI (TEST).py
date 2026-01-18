# ======================================================================================
# CAPEX AI RT2026 
#
#
# requirements.txt:
# streamlit
# pandas
# numpy
# scipy
# scikit-learn
# plotly
# matplotlib
# python-pptx
# openpyxl
# requests
# ======================================================================================

import io
import json
import zipfile
import re
import requests
import numpy as np
import pandas as pd
import streamlit as st

# --- Guard: scikit-learn missing (prevents hard crash on Streamlit Cloud) ---
try:
    from sklearn.impute import KNNImputer, SimpleImputer
    from sklearn.model_selection import train_test_split
    from sklearn.preprocessing import MinMaxScaler
    from sklearn.ensemble import RandomForestRegressor, GradientBoostingRegressor
    from sklearn.linear_model import Ridge, Lasso
    from sklearn.svm import SVR
    from sklearn.tree import DecisionTreeRegressor
    from sklearn.pipeline import Pipeline
    from sklearn.metrics import mean_squared_error, r2_score
except Exception as e:
    st.error(
        "❌ Missing dependency: **scikit-learn**.\n\n"
        "Fix:\n"
        "1) Open your **requirements.txt**\n"
        "2) Add this line: `scikit-learn`\n"
        "3) Commit + redeploy.\n\n"
        f"Details: {e}"
    )
    st.stop()

from scipy.stats import linregress

import plotly.express as px
import plotly.graph_objects as go

import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from openpyxl.formatting.rule import ColorScaleRule
from openpyxl.chart import BarChart, LineChart, Reference
from openpyxl.utils import get_column_letter

# ---------------------------------------------------------------------------------------
# PAGE CONFIG
# ---------------------------------------------------------------------------------------
st.set_page_config(
    page_title="CAPEX AI RT2026",
    page_icon="💠",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ---------------------------------------------------------------------------------------
# THEME TOKENS
# ---------------------------------------------------------------------------------------
PETRONAS = {
    "teal": "#00A19B",
    "teal_dark": "#008C87",
    "purple": "#6C4DD3",
    "white": "#FFFFFF",
    "black": "#0E1116",
    "border": "rgba(0,0,0,0.10)",
}

# ---------------------------------------------------------------------------------------
# SHAREPOINT LINKS (FILL THESE LATER)
# ---------------------------------------------------------------------------------------
SHAREPOINT_LINKS = {
    "Shallow Water": "https://petronas.sharepoint.com/sites/your-site/shallow-water",
    "Deep Water": "https://petronas.sharepoint.com/sites/your-site/deep-water",
    "Onshore": "https://petronas.sharepoint.com/sites/your-site/onshore",
    "Uncon": "https://petronas.sharepoint.com/sites/your-site/uncon",
    "CCS": "https://petronas.sharepoint.com/sites/your-site/ccs",
}

# ---------------------------------------------------------------------------------------
# GLOBAL CSS
# ---------------------------------------------------------------------------------------
st.markdown(
    f"""
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700;800&display=swap');
html, body {{ font-family: 'Inter', sans-serif; }}
[data-testid="stAppViewContainer"] {{
  background: {PETRONAS["white"]};
  color: {PETRONAS["black"]};
  padding-top: 0.5rem;
}}
#MainMenu, footer {{ visibility: hidden; }}
[data-testid="stSidebar"] {{
  background: linear-gradient(180deg, {PETRONAS["teal"]} 0%, {PETRONAS["teal_dark"]} 100%) !important;
  color: #fff !important;
  border-top-right-radius: 16px;
  border-bottom-right-radius: 16px;
  box-shadow: 0 6px 20px rgba(0,0,0,0.15);
}}
[data-testid="stSidebar"] * {{ color: #fff !important; }}
[data-testid="collapsedControl"] {{
  position: fixed !important;
  top: 50% !important;
  left: 10px !important;
  transform: translateY(-50%) !important;
  z-index: 9999 !important;
}}
.petronas-hero {{
  border-radius: 20px;
  padding: 28px 32px;
  margin: 6px 0 18px 0;
  color: #fff;
  background: linear-gradient(135deg, {PETRONAS["teal"]}, {PETRONAS["purple"]}, {PETRONAS["black"]});
  background-size: 200% 200%;
  animation: heroGradient 8s ease-in-out infinite, fadeIn .8s ease-in-out, heroPulse 5s ease-in-out infinite;
  box-shadow: 0 10px 24px rgba(0,0,0,.12);
}}
@keyframes heroGradient {{
  0% {{ background-position: 0% 50%; }}
  50% {{ background-position: 100% 50%; }}
  100% {{ background-position: 0% 50%; }}
}}
@keyframes fadeIn {{
  from {{ opacity: 0; transform: translateY(10px); }}
  to {{ opacity: 1; transform: translateY(0); }}
}}
@keyframes heroPulse {{
  0%   {{ box-shadow: 0 0 16px rgba(0,161,155,0.45); }}
  25%  {{ box-shadow: 0 0 26px rgba(108,77,211,0.55); }}
  50%  {{ box-shadow: 0 0 36px rgba(0,161,155,0.55); }}
  75%  {{ box-shadow: 0 0 26px rgba(108,77,211,0.55); }}
  100% {{ box-shadow: 0 0 16px rgba(0,161,155,0.45); }}
}}
.petronas-hero h1 {{ margin: 0 0 5px; font-weight: 800; letter-spacing: 0.3px; }}
.petronas-hero p {{ margin: 0; opacity: .9; font-weight: 500; }}

.stButton > button, .stDownloadButton > button, .petronas-button {{
  border-radius: 10px;
  padding: .6rem 1.1rem;
  font-weight: 600;
  color: #fff !important;
  border: none;
  background: linear-gradient(to right, {PETRONAS["teal"]}, {PETRONAS["purple"]});
  background-size: 200% auto;
  transition: background-position .85s ease, transform .2s ease, box-shadow .25s ease;
  text-decoration: none;
  display: inline-block;
}}
.stButton > button:hover, .stDownloadButton > button:hover, .petronas-button:hover {{
  background-position: right center;
  transform: translateY(-1px);
  box-shadow: 0 6px 16px rgba(0,0,0,0.18);
}}

.stTabs [role="tablist"] {{
  display: flex;
  gap: 8px;
  border-bottom: none;
  padding-bottom: 6px;
}}
.stTabs [role="tab"] {{
  background: #fff;
  color: {PETRONAS["black"]};
  border-radius: 8px;
  padding: 10px 18px;
  border: 1px solid {PETRONAS["border"]};
  font-weight: 600;
  transition: all .3s ease;
  position: relative;
}}
.stTabs [role="tab"]:hover {{
  background: linear-gradient(to right, {PETRONAS["teal"]}, {PETRONAS["purple"]});
  color: #fff;
}}
.stTabs [role="tab"][aria-selected="true"] {{
  background: linear-gradient(to right, {PETRONAS["teal"]}, {PETRONAS["purple"]});
  color: #fff;
  border-color: transparent;
  box-shadow: 0 4px 16px rgba(0,0,0,0.15);
}}
.stTabs [role="tab"][aria-selected="true"]::after {{
  content: "";
  position: absolute;
  left: 10%;
  bottom: -3px;
  width: 80%;
  height: 3px;
  background: linear-gradient(90deg, {PETRONAS["teal"]}, {PETRONAS["purple"]}, {PETRONAS["teal"]});
  background-size: 200% 100%;
  border-radius: 2px;
  animation: glowSlide 2.5s linear infinite;
}}
@keyframes glowSlide {{
  0% {{ background-position: 0% 50%; }}
  50% {{ background-position: 100% 50%; }}
  100% {{ background-position: 0% 50%; }}
}}
</style>
""",
    unsafe_allow_html=True,
)

# ---------------------------------------------------------------------------------------
# HERO HEADER
# ---------------------------------------------------------------------------------------
st.markdown(
    """
<div class="petronas-hero">
  <h1>CAPEX AI RT2026</h1>
  <p>Data-driven CAPEX prediction</p>
</div>
""",
    unsafe_allow_html=True,
)

# ---------------------------------------------------------------------------------------
# AUTH
# ---------------------------------------------------------------------------------------
if "authenticated" not in st.session_state:
    st.session_state.authenticated = False

APPROVED_EMAILS = [str(e).strip().lower() for e in st.secrets.get("emails", [])]
correct_password = st.secrets.get("password", None)

if not st.session_state.authenticated:
    with st.form("login_form"):
        st.markdown("#### 🔐 Access Required", unsafe_allow_html=True)
        email = (st.text_input("Email Address", key="login_email") or "").strip().lower()
        password = st.text_input("Access Password", type="password", key="login_pwd")
        submitted = st.form_submit_button("Login")

        if submitted:
            ok = (email in APPROVED_EMAILS) and (password == correct_password)
            if ok:
                st.session_state.authenticated = True
                st.success("✅ Access granted.")
                st.rerun()
            else:
                st.error("❌ Invalid credentials.")
    st.stop()

# ---------------------------------------------------------------------------------------
# SESSION STATE
# ---------------------------------------------------------------------------------------
if "datasets" not in st.session_state:
    st.session_state.datasets = {}
if "predictions" not in st.session_state:
    st.session_state.predictions = {}
if "processed_excel_files" not in st.session_state:
    st.session_state.processed_excel_files = set()
if "_last_metrics" not in st.session_state:
    st.session_state._last_metrics = None
if "projects" not in st.session_state:
    st.session_state.projects = {}
if "component_labels" not in st.session_state:
    st.session_state.component_labels = {}
if "uploader_nonce" not in st.session_state:
    st.session_state.uploader_nonce = 0
if "widget_nonce" not in st.session_state:
    st.session_state.widget_nonce = 0

# ---------------------------------------------------------------------------------------
# HELPERS
# ---------------------------------------------------------------------------------------
def toast(msg, icon="✅"):
    try:
        st.toast(f"{icon} {msg}")
    except Exception:
        st.success(msg if icon == "✅" else msg)


def format_with_commas(num):
    try:
        return f"{float(num):,.2f}"
    except Exception:
        return str(num)


def normalize_to_100(d: dict):
    total = sum(float(v) for v in d.values())
    if total <= 0:
        return d, total
    out = {k: float(v) * 100.0 / total for k, v in d.items()}
    keys = list(out.keys())
    rounded = {k: round(out[k], 2) for k in keys}
    diff = 100.0 - sum(rounded.values())
    if keys:
        rounded[keys[-1]] = round(rounded[keys[-1]] + diff, 2)
    return rounded, total


def is_junk_col(colname: str) -> bool:
    h = str(colname).strip().upper()
    return (not h) or h.startswith("UNNAMED") or h in {"INDEX", "IDX"}


def currency_from_header(header: str) -> str:
    h = (header or "").strip().upper()
    if "€" in h:
        return "€"
    if "£" in h:
        return "£"
    if "$" in h:
        return "$"
    if re.search(r"\bUSD\b", h):
        return "USD"
    if re.search(r"\b(MYR|RM)\b", h):
        return "RM"
    return ""


def get_currency_symbol(df: pd.DataFrame, target_col: str | None = None) -> str:
    if df is None or df.empty:
        return ""
    if target_col and target_col in df.columns:
        return currency_from_header(str(target_col))
    for c in reversed(df.columns):
        if not is_junk_col(c):
            return currency_from_header(str(c))
    return ""


def cost_breakdown(
    base_pred: float,
    eprr: dict,
    sst_pct: float,
    owners_pct: float,
    cont_pct: float,
    esc_pct: float,
):
    base_pred = float(base_pred)

    owners_cost = round(base_pred * (owners_pct / 100.0), 2)
    sst_cost = round(base_pred * (sst_pct / 100.0), 2)

    contingency_cost = round((base_pred + owners_cost) * (cont_pct / 100.0), 2)
    escalation_cost = round((base_pred + owners_cost) * (esc_pct / 100.0), 2)

    eprr_costs = {k: round(base_pred * (float(v) / 100.0), 2) for k, v in (eprr or {}).items()}

    # FIX: includes SST
    grand_total = round(base_pred + owners_cost + sst_cost + contingency_cost + escalation_cost, 2)
    return owners_cost, sst_cost, contingency_cost, escalation_cost, eprr_costs, grand_total


def project_components_df(proj):
    comps = proj.get("components", [])
    rows = []
    for c in comps:
        rows.append(
            {
                "Component": c["component_type"],
                "Dataset": c["dataset"],
                "Model": c.get("model_used", ""),
                "Base CAPEX": float(c["prediction"]),
                "Owner's Cost": float(c["breakdown"]["owners_cost"]),
                "Contingency": float(c["breakdown"]["contingency_cost"]),
                "Escalation": float(c["breakdown"]["escalation_cost"]),
                "SST": float(c["breakdown"]["sst_cost"]),
                "Grand Total": float(c["breakdown"]["grand_total"]),
            }
        )
    return pd.DataFrame(rows)


def project_totals(proj):
    dfc = project_components_df(proj)
    if dfc.empty:
        return {"capex_sum": 0.0, "owners": 0.0, "cont": 0.0, "esc": 0.0, "sst": 0.0, "grand_total": 0.0}
    return {
        "capex_sum": float(dfc["Base CAPEX"].sum()),
        "owners": float(dfc["Owner's Cost"].sum()),
        "cont": float(dfc["Contingency"].sum()),
        "esc": float(dfc["Escalation"].sum()),
        "sst": float(dfc["SST"].sum()),
        "grand_total": float(dfc["Grand Total"].sum()),
    }


# ======================================================================================
# TARGET ALWAYS LAST COLUMN (with safe numeric coercion)
# ======================================================================================
def get_last_column_target(df: pd.DataFrame) -> str:
    if df is None or df.empty:
        raise ValueError("Empty dataset.")
    return str(df.columns[-1])


def coerce_series_numeric(s: pd.Series) -> pd.Series:
    if pd.api.types.is_numeric_dtype(s):
        return s.astype(float)
    return pd.to_numeric(s, errors="coerce")


def numeric_features_from_df(df: pd.DataFrame, target_col: str) -> tuple[pd.DataFrame, pd.Series]:
    """
    - Target: ALWAYS df[target_col] (last column by default)
    - Coerce target to numeric (errors->NaN)
    - Features: numeric columns excluding target (coerce only numeric dtypes)
    - If target becomes all NaN -> fallback to last numeric column (warn)
    """
    if target_col not in df.columns:
        raise ValueError(f"Target column not found: {target_col}")

    y_raw = df[target_col]
    y = coerce_series_numeric(y_raw)

    # numeric feature selection (exclude target)
    X = df.drop(columns=[target_col]).select_dtypes(include=[np.number]).copy()

    # if no numeric feature columns, try to coerce other columns? (skip to keep model sane)
    if X.shape[1] < 1:
        raise ValueError("Need at least 1 numeric feature column (excluding target).")

    # If target is unusable, fallback to last numeric column
    if y.dropna().shape[0] == 0:
        num_cols = df.select_dtypes(include=[np.number]).columns.tolist()
        if len(num_cols) >= 2:
            fallback = num_cols[-1]
            st.warning(
                f"⚠️ Last column '{target_col}' could not be converted to numeric. "
                f"Falling back to last numeric column '{fallback}'."
            )
            y = df[fallback].astype(float)
            X = df.drop(columns=[fallback]).select_dtypes(include=[np.number]).copy()
            target_col = str(fallback)
        else:
            raise ValueError(
                f"Last column '{target_col}' is not numeric and no numeric fallback target exists."
            )

    return X, y.astype(float)


# ======================================================================================
# ✅ MONTE CARLO HELPERS
# ======================================================================================
def _coerce_float(x, default=np.nan):
    try:
        if x is None:
            return default
        if isinstance(x, str) and x.strip() == "":
            return default
        return float(x)
    except Exception:
        return default


def monte_carlo_component(
    model_pipe: Pipeline,
    feature_cols: list[str],
    base_payload: dict,
    n_sims: int = 5000,
    seed: int = 42,
    feature_sigma_pct: float = 5.0,
    pct_sigma_abs: float = 1.0,
    eprr: dict | None = None,
    sst_pct: float = 0.0,
    owners_pct: float = 0.0,
    cont_pct: float = 0.0,
    esc_pct: float = 0.0,
    normalize_eprr_each_draw: bool = False,
) -> pd.DataFrame:
    rng = np.random.default_rng(int(seed))
    n = int(n_sims)

    base_vec = np.array([_coerce_float(base_payload.get(c), np.nan) for c in feature_cols], dtype=float)
    Xsim = np.tile(base_vec, (n, 1))

    # multiplicative feature noise on non-NaN
    sigma = float(feature_sigma_pct) / 100.0
    if sigma > 0:
        noise = rng.normal(0.0, sigma, size=Xsim.shape)
        mask = ~np.isnan(Xsim)
        Xsim[mask] = Xsim[mask] * (1.0 + noise[mask])

    df_sim = pd.DataFrame(Xsim, columns=feature_cols)
    base_preds = model_pipe.predict(df_sim).astype(float)

    p_sig = float(pct_sigma_abs)
    sst_draw = np.clip(rng.normal(loc=float(sst_pct), scale=p_sig, size=n), 0.0, 100.0)
    own_draw = np.clip(rng.normal(loc=float(owners_pct), scale=p_sig, size=n), 0.0, 100.0)
    con_draw = np.clip(rng.normal(loc=float(cont_pct), scale=p_sig, size=n), 0.0, 100.0)
    esc_draw = np.clip(rng.normal(loc=float(esc_pct), scale=p_sig, size=n), 0.0, 100.0)

    eprr = eprr or {}
    e_keys = list(eprr.keys())
    e_mat = None
    if e_keys:
        e_mat = np.vstack(
            [np.clip(rng.normal(loc=float(eprr.get(k, 0.0)), scale=p_sig, size=n), 0.0, 100.0) for k in e_keys]
        ).T
        if normalize_eprr_each_draw:
            rs = e_mat.sum(axis=1)
            rs[rs == 0] = 1.0
            e_mat = (e_mat / rs[:, None]) * 100.0

    owners_cost = np.round(base_preds * (own_draw / 100.0), 2)
    sst_cost = np.round(base_preds * (sst_draw / 100.0), 2)
    contingency_cost = np.round((base_preds + owners_cost) * (con_draw / 100.0), 2)
    escalation_cost = np.round((base_preds + owners_cost) * (esc_draw / 100.0), 2)
    grand_total = np.round(base_preds + owners_cost + sst_cost + contingency_cost + escalation_cost, 2)

    out = pd.DataFrame(
        {
            "base_pred": base_preds,
            "owners_cost": owners_cost,
            "sst_cost": sst_cost,
            "contingency_cost": contingency_cost,
            "escalation_cost": escalation_cost,
            "grand_total": grand_total,
            "owners_pct": own_draw,
            "sst_pct": sst_draw,
            "cont_pct": con_draw,
            "esc_pct": esc_draw,
        }
    )

    if e_keys:
        for j, k in enumerate(e_keys):
            out[f"eprr_{k}_pct"] = e_mat[:, j]
            out[f"eprr_{k}_cost"] = np.round(base_preds * (e_mat[:, j] / 100.0), 2)

    return out


def scenario_bucket_from_baseline(values: pd.Series, baseline: float, low_cut_pct: float, band_pct: float, high_cut_pct: float):
    baseline = float(baseline) if np.isfinite(baseline) and float(baseline) != 0 else float(values.median())
    pct_delta = (values - baseline) / baseline

    def label(v):
        if v < (-low_cut_pct / 100.0):
            return "Low"
        if (-band_pct / 100.0) <= v <= (band_pct / 100.0):
            return "Base"
        if v > (high_cut_pct / 100.0):
            return "High"
        return "Unbucketed"

    buckets = pct_delta.apply(label)
    return buckets, pct_delta * 100.0


# ---------------------------------------------------------------------------------------
# DATA / MODEL HELPERS
# ---------------------------------------------------------------------------------------
GITHUB_USER = "Hafizuddin-Abd-Rahman-Dev-Upstream"
REPO_NAME = "Cost-Predictor"
BRANCH = "main"
DATA_FOLDER = "pages/data_CAPEX"

MODEL_CANDIDATES = {
    "RandomForest": lambda rs=42: RandomForestRegressor(random_state=rs),
    "GradientBoosting": lambda rs=42: GradientBoostingRegressor(random_state=rs),
    "Ridge": lambda rs=42: Ridge(),
    "Lasso": lambda rs=42: Lasso(),
    "SVR": lambda rs=42: SVR(),
    "DecisionTree": lambda rs=42: DecisionTreeRegressor(random_state=rs),
}
SCALE_MODELS = {"Ridge", "Lasso", "SVR"}


@st.cache_data(ttl=600, show_spinner=False)
def fetch_json(url: str):
    r = requests.get(url, timeout=15)
    r.raise_for_status()
    return r.json()


@st.cache_data(ttl=600, show_spinner=False)
def list_csvs_from_manifest(folder_path: str):
    manifest_url = f"https://raw.githubusercontent.com/{GITHUB_USER}/{REPO_NAME}/{BRANCH}/{folder_path}/files.json"
    try:
        data = fetch_json(manifest_url)
        if isinstance(data, list):
            return [str(x) for x in data]
        return []
    except Exception as e:
        st.error(f"Failed to load CSV manifest: {e}")
        return []


def make_pipeline(model_name: str, random_state=42):
    ctor = MODEL_CANDIDATES[model_name]
    try:
        model = ctor(random_state)
    except TypeError:
        model = ctor()

    steps = [("imputer", SimpleImputer(strategy="median"))]
    if model_name in SCALE_MODELS:
        steps.append(("scaler", MinMaxScaler()))
    steps.append(("model", model))
    return Pipeline(steps)


def evaluate_models(X, y, test_size=0.2, random_state=42):
    Xtr, Xte, ytr, yte = train_test_split(X, y, test_size=test_size, random_state=random_state)

    rows = []
    best_name = None
    best_r2 = -np.inf
    best_rmse = None

    for name in MODEL_CANDIDATES.keys():
        pipe = make_pipeline(name, random_state=random_state)
        pipe.fit(Xtr, ytr)
        yhat = pipe.predict(Xte)
        rmse = float(np.sqrt(mean_squared_error(yte, yhat)))
        r2 = float(r2_score(yte, yhat))
        rows.append({"model": name, "rmse": rmse, "r2": r2})
        if r2 > best_r2:
            best_r2 = r2
            best_rmse = rmse
            best_name = name

    rows_sorted = sorted(rows, key=lambda d: d["r2"], reverse=True)
    metrics = {"best_model": best_name, "rmse": best_rmse, "r2": best_r2, "models": rows_sorted}
    return metrics


@st.cache_resource(show_spinner=False)
def train_best_model_cached(df: pd.DataFrame, target_col: str, test_size: float, random_state: int, dataset_key: str):
    X, y = numeric_features_from_df(df, target_col)
    metrics = evaluate_models(X, y, test_size=test_size, random_state=random_state)
    best_name = metrics.get("best_model") or "RandomForest"
    best_pipe = make_pipeline(best_name, random_state=random_state)
    best_pipe.fit(X, y)
    return best_pipe, metrics, list(X.columns), target_col, best_name


def single_prediction(model_pipe: Pipeline, feature_cols: list[str], payload: dict):
    row = {}
    for c in feature_cols:
        v = payload.get(c, np.nan)
        try:
            if v is None or (isinstance(v, str) and v.strip() == ""):
                row[c] = np.nan
            else:
                row[c] = float(v)
        except Exception:
            row[c] = np.nan
    df_in = pd.DataFrame([row], columns=feature_cols)
    return float(model_pipe.predict(df_in)[0])


@st.cache_data(show_spinner=False, ttl=600)
def knn_impute_numeric(df: pd.DataFrame, k: int = 5):
    # purely numeric
    num = df.select_dtypes(include=[np.number]).copy()
    num = num.dropna(axis=1, how="all")
    if num.shape[1] < 2:
        raise ValueError("Need at least 2 numeric columns for KNN viz.")
    arr = KNNImputer(n_neighbors=k).fit_transform(num)
    return pd.DataFrame(arr, columns=num.columns)


# ---------------------------------------------------------------------------------------
# NAV ROW — SHAREPOINT BUTTONS
# ---------------------------------------------------------------------------------------
nav_labels = ["SHALLOW WATER", "DEEP WATER", "ONSHORE", "UNCON", "CCS"]
nav_cols = st.columns(len(nav_labels))
for col, label in zip(nav_cols, nav_labels):
    with col:
        url = SHAREPOINT_LINKS.get(label.title(), "#")
        st.markdown(
            f'''
            <a href="{url}" target="_blank" rel="noopener"
               class="petronas-button"
               style="width:100%; text-align:center; display:inline-block;">
               {label}
            </a>
            ''',
            unsafe_allow_html=True,
        )

# ---------------------------------------------------------------------------------------
# TOP-LEVEL TABS
# ---------------------------------------------------------------------------------------
tab_data, tab_pb, tab_mc, tab_compare = st.tabs(
    ["📊 Data", "🏗️ Project Builder", "🎲 Monte Carlo", "🔀 Compare Projects"]
)

# =======================================================================================
# DATA TAB
# =======================================================================================
with tab_data:
    st.markdown('<h3 style="margin-top:0;color:#000;">📁 Data</h3>', unsafe_allow_html=True)

    st.markdown('<h4 style="margin:0;color:#000;">Data Sources</h4><p></p>', unsafe_allow_html=True)
    c1, c2 = st.columns([1.2, 1])
    with c1:
        data_source = st.radio("Choose data source", ["Upload CSV", "Load from Server"], horizontal=True, key="data_source")

    with c2:
        st.caption("Enterprise Storage (SharePoint)")
        data_link = (
            "https://petronas.sharepoint.com/sites/ecm_ups_coe/confidential/"
            "DFE%20Cost%20Engineering/Forms/AllItems.aspx?"
            "id=%2Fsites%2Fecm%5Fups%5Fcoe%2Fconfidential%2FDFE%20Cost%20Engineering"
            "%2F2%2ETemplate%20Tools%2FCost%20Predictor%2FDatabase%2FCAPEX%20%2D%20RT%20Q1%202025"
        )
        st.markdown(
            f'<a href="{data_link}" target="_blank" rel="noopener" class="petronas-button">Open Enterprise Storage</a>',
            unsafe_allow_html=True,
        )

    uploaded_files = []
    if data_source == "Upload CSV":
        uploaded_files = st.file_uploader(
            "Upload CSV files (max 200MB)",
            type="csv",
            accept_multiple_files=True,
            key=f"csv_uploader_{st.session_state.uploader_nonce}",
        )
    else:
        github_csvs = list_csvs_from_manifest(DATA_FOLDER)
        if github_csvs:
            selected_file = st.selectbox("Choose CSV from GitHub", github_csvs, key="github_csv_select")
            if st.button("Load selected CSV", key="load_github_csv_btn"):
                raw_url = f"https://raw.githubusercontent.com/{GITHUB_USER}/{REPO_NAME}/{BRANCH}/{DATA_FOLDER}/{selected_file}"
                try:
                    df = pd.read_csv(raw_url)
                    st.session_state.datasets[selected_file] = df
                    st.session_state.predictions.setdefault(selected_file, [])
                    toast(f"Loaded from GitHub: {selected_file}")
                    st.rerun()
                except Exception as e:
                    st.error(f"Error loading CSV: {e}")
        else:
            st.info("No CSV files found in GitHub folder.")

    if uploaded_files:
        for up in uploaded_files:
            if up.name not in st.session_state.datasets:
                try:
                    df = pd.read_csv(up)
                    st.session_state.datasets[up.name] = df
                    st.session_state.predictions.setdefault(up.name, [])
                except Exception as e:
                    st.error(f"Failed to read {up.name}: {e}")
        toast("Dataset(s) added.")

    st.divider()

    cA, cB, cC, cD = st.columns([1, 1, 1, 2])
    with cA:
        if st.button("🧹 Clear all predictions", key="clear_preds_btn"):
            st.session_state.predictions = {k: [] for k in st.session_state.predictions.keys()}
            toast("All predictions cleared.", "🧹")
            st.rerun()

    with cB:
        if st.button("🧺 Clear processed files history", key="clear_processed_btn"):
            st.session_state.processed_excel_files = set()
            toast("Processed files history cleared.", "🧺")
            st.rerun()

    with cC:
        if st.button("🔁 Refresh server manifest", key="refresh_manifest_btn"):
            list_csvs_from_manifest.clear()
            fetch_json.clear()
            toast("Server manifest refreshed.", "🔁")
            st.rerun()

    with cD:
        if st.button("🗂️ Clear all uploaded / loaded files (keep projects)", key="clear_datasets_btn"):
            st.session_state.datasets = {}
            st.session_state.predictions = {}
            st.session_state.processed_excel_files = set()
            st.session_state._last_metrics = None
            st.session_state.uploader_nonce += 1
            st.session_state.widget_nonce += 1
            toast("All datasets cleared. Projects preserved.", "🗂️")
            st.rerun()

    st.divider()

    # -------------------------
    # Active dataset preview + target selection (NOW: always last column)
    # -------------------------
    if st.session_state.datasets:
        ds_name_data = st.selectbox("Active dataset", list(st.session_state.datasets.keys()), key="active_dataset_data")
        df_active = st.session_state.datasets[ds_name_data]

        target_col_active = get_last_column_target(df_active)
        st.session_state[f"target_col__{ds_name_data}"] = target_col_active

        currency_active = get_currency_symbol(df_active, target_col_active)

        colA, colB, colC, colD2 = st.columns([1, 1, 1, 2])
        with colA:
            st.metric("Rows", f"{df_active.shape[0]:,}")
        with colB:
            st.metric("Columns", f"{df_active.shape[1]:,}")
        with colC:
            st.metric("Currency", f"{currency_active or '—'}")
        with colD2:
            st.caption(f"Target is forced to LAST column: **{target_col_active}**")

        with st.expander("Preview (first 10 rows)", expanded=False):
            st.dataframe(df_active.head(10), use_container_width=True)
    else:
        st.info("Upload or load a dataset to proceed.")
        st.stop()

    # ========================= MODEL TRAINING =================================
    st.divider()
    st.markdown('<h3 style="margin-top:0;color:#000;">⚙️ Model</h3>', unsafe_allow_html=True)

    ds_name_model = st.selectbox("Dataset for model training", list(st.session_state.datasets.keys()), key="ds_model")
    df_model = st.session_state.datasets[ds_name_model]
    target_col_model = get_last_column_target(df_model)
    st.session_state[f"target_col__{ds_name_model}"] = target_col_model

    m1, m2 = st.columns([1, 3])
    with m1:
        test_size = st.slider("Test size", 0.1, 0.5, 0.2, 0.05, key="train_test_size")
        run_train = st.button("Run training", key="run_training_btn")
    with m2:
        st.caption(f"Best-model selection over 6 regressors. Target forced to last column: {target_col_model}")

    if run_train:
        try:
            with st.spinner("Training model (cached)..."):
                pipe, metrics, feat_cols, y_name, best_name = train_best_model_cached(
                    df_model,
                    target_col_model,
                    test_size=float(test_size),
                    random_state=42,
                    dataset_key=ds_name_model,
                )
            st.session_state._last_metrics = metrics
            st.session_state[f"best_model__{ds_name_model}"] = best_name
            toast("Training complete.")
            c_rmse, c_r2, c_best = st.columns(3)
            with c_rmse:
                st.metric("RMSE (best)", f"{metrics['rmse']:,.2f}")
            with c_r2:
                st.metric("R² (best)", f"{metrics['r2']:.3f}")
            with c_best:
                st.metric("Best Model", best_name)

            models_list = metrics.get("models", [])
            if models_list:
                df_models = pd.DataFrame(models_list).set_index("model")
                st.markdown("##### Model comparison")
                st.dataframe(df_models, use_container_width=True)
        except Exception as e:
            st.error(f"Training failed: {e}")

    # ========================= VISUALIZATION ==================================
    st.divider()
    st.markdown('<h3 style="margin-top:0;color:#000;">📈 Visualization</h3>', unsafe_allow_html=True)

    ds_name_viz = st.selectbox("Dataset for visualization", list(st.session_state.datasets.keys()), key="ds_viz")
    df_viz = st.session_state.datasets[ds_name_viz]
    target_col_viz = get_last_column_target(df_viz)
    st.session_state[f"target_col__{ds_name_viz}"] = target_col_viz

    with st.expander("Visualization settings", expanded=False):
        use_knn = st.checkbox("Use KNN imputation for visualization (slower)", value=False, key="viz_knn")
        knn_k = st.slider("KNN neighbors", 2, 15, 5, 1, disabled=not use_knn, key="viz_knn_k")

    try:
        # Viz: numeric-only matrix (best effort)
        if use_knn:
            num_imputed = knn_impute_numeric(df_viz, k=int(knn_k))
        else:
            num = df_viz.select_dtypes(include=[np.number]).copy()
            num = num.dropna(axis=1, how="all")
            if num.shape[1] < 2:
                raise ValueError("Not enough numeric columns for visualization.")
            num_imputed = pd.DataFrame(SimpleImputer(strategy="median").fit_transform(num), columns=num.columns)

        st.markdown('<h4 style="margin:0;color:#000;">Correlation Matrix</h4><p>Exploration</p>', unsafe_allow_html=True)
        corr = num_imputed.corr(numeric_only=True)
        fig_corr = px.imshow(corr, text_auto=".2f", aspect="auto", color_continuous_scale="RdBu_r", zmin=-1, zmax=1)
        fig_corr.update_layout(margin=dict(l=0, r=0, t=10, b=0))
        st.plotly_chart(fig_corr, use_container_width=True)

        st.markdown('<h4 style="margin:0;color:#000;">Feature Importance (RandomForest)</h4><p>Model insight</p>', unsafe_allow_html=True)

        # choose a target for viz if last column not numeric; use last numeric
        viz_target = num_imputed.columns[-1]
        X_viz = num_imputed.drop(columns=[viz_target])
        y_viz = num_imputed[viz_target]

        model_viz = RandomForestRegressor(random_state=42).fit(X_viz, y_viz)
        importances = model_viz.feature_importances_
        fi = pd.DataFrame({"feature": X_viz.columns, "importance": importances}).sort_values("importance", ascending=True)
        fig_fi = go.Figure(go.Bar(x=fi["importance"], y=fi["feature"], orientation="h"))
        fig_fi.update_layout(xaxis_title="Importance", yaxis_title="Feature", margin=dict(l=0, r=0, t=10, b=0))
        st.plotly_chart(fig_fi, use_container_width=True)

        st.markdown('<h4 style="margin:0;color:#000;">Cost Curve</h4><p>Trend</p>', unsafe_allow_html=True)
        feat = st.selectbox("Select feature for cost curve", list(X_viz.columns), key="viz_cost_curve_feat")
        x_vals = X_viz[feat].to_numpy(dtype=float)
        y_vals = y_viz.to_numpy(dtype=float)
        mask = (~np.isnan(x_vals)) & (~np.isnan(y_vals))
        scatter_df = pd.DataFrame({feat: x_vals[mask], "Target": y_vals[mask]})
        fig_cc = px.scatter(scatter_df, x=feat, y="Target", opacity=0.65)

        if mask.sum() >= 2 and np.unique(x_vals[mask]).size >= 2:
            xv = scatter_df[feat].to_numpy(dtype=float)
            yv = scatter_df["Target"].to_numpy(dtype=float)
            slope, intercept, r_value, p_value, std_err = linregress(xv, yv)
            x_line = np.linspace(xv.min(), xv.max(), 100)
            y_line = slope * x_line + intercept
            fig_cc.add_trace(
                go.Scatter(
                    x=x_line,
                    y=y_line,
                    mode="lines",
                    name=f"Fit: y={slope:.2f}x+{intercept:.2f} (R²={r_value**2:.3f})",
                )
            )
        else:
            st.warning("Not enough valid/variable data to compute regression.")
        fig_cc.update_layout(margin=dict(l=0, r=0, t=10, b=0))
        st.plotly_chart(fig_cc, use_container_width=True)

    except Exception as e:
        st.error(f"Visualization error: {e}")

    # ========================= PREDICT =======================================
    st.divider()
    st.markdown('<h3 style="margin-top:0;color:#000;">🎯 Predict</h3>', unsafe_allow_html=True)

    ds_name_pred = st.selectbox("Dataset for prediction", list(st.session_state.datasets.keys()), key="ds_pred")
    df_pred = st.session_state.datasets[ds_name_pred]
    target_col_pred = get_last_column_target(df_pred)
    st.session_state[f"target_col__{ds_name_pred}"] = target_col_pred
    currency_pred = get_currency_symbol(df_pred, target_col_pred)

    st.markdown('<h4 style="margin:0;color:#000;">WBS Level 1</h4><p>Step 1</p>', unsafe_allow_html=True)
    c1, c2 = st.columns([1, 1])
    with c1:
        st.markdown("(use +/-)")
        eng = st.number_input("Engineering (%)", min_value=0.0, max_value=100.0, value=0.0, step=1.0, key="pred_eng")  # Changed to 0.0
        procurement = st.number_input("Procurement (%)", min_value=0.0, max_value=100.0, value=0.0, step=1.0, key="pred_procurement")  # Changed to 0.0
        fabrication = st.number_input("Fabrication/Construction (%)", min_value=0.0, max_value=100.0, value=0.0, step=1.0, key="pred_fabrication")  # Changed to 0.0
        ti = st.number_input("Transportation & Installation (T&I) (%)", min_value=0.0, max_value=100.0, value=0.0, step=1.0, key="pred_ti")  # Changed to 0.0

        eprr = {"Engineering": eng, "Procurement": procurement, "Fabrication/Construction": fabrication, "Transportation & Installation": ti}
        eprr_total = sum(eprr.values())
        st.caption(f"WBS total: **{eprr_total:.2f}%**")

    with c2:
        st.markdown("(use +/-)")
        sst_pct = st.number_input("SST (%)", min_value=0.0, max_value=100.0, value=0.0, step=0.5, key="pred_sst")
        owners_pct = st.number_input("Owner's Cost (%)", min_value=0.0, max_value=100.0, value=0.0, step=0.5, key="pred_owner")
        cont_pct = st.number_input("Contingency (%)", min_value=0.0, max_value=100.0, value=0.0, step=0.5, key="pred_cont")
        esc_pct = st.number_input("Escalation & Inflation (%)", min_value=0.0, max_value=100.0, value=0.0, step=0.5, key="pred_esc")

    st.markdown('<h4 style="margin:0;color:#000;">Predict (Single)</h4><p>Step 2</p>', unsafe_allow_html=True)
    project_name = st.text_input("Project Name", placeholder="e.g., Offshore Pipeline Replacement 2026", key="pred_project_name")

    try:
        pipe, metrics_auto, feat_cols, y_name, best_name = train_best_model_cached(
            df_pred,
            target_col_pred,
            test_size=0.2,
            random_state=42,
            dataset_key=ds_name_pred,
        )
        st.session_state._last_metrics = st.session_state._last_metrics or metrics_auto
        st.caption(f"Using cached best model: **{best_name}** | Target (last col): **{y_name}**")
    except Exception as e:
        st.error(f"Model setup failed: {e}")
        st.stop()

    st.caption("Provide feature values (1 row). Leave blank for NaN.")
    input_key = f"input_row__{ds_name_pred}"
    if input_key not in st.session_state:
        st.session_state[input_key] = {c: np.nan for c in feat_cols}

    row_df = pd.DataFrame([st.session_state[input_key]], columns=feat_cols)
    pred_editor_key = f"pred_editor__{ds_name_pred}__{st.session_state.widget_nonce}"
    edited = st.data_editor(row_df, num_rows="fixed", use_container_width=True, key=pred_editor_key)
    payload = edited.iloc[0].to_dict()

    if st.button("Run Prediction", key="run_pred_btn"):
        try:
            pred_val = single_prediction(pipe, feat_cols, payload)

            owners_cost, sst_cost, contingency_cost, escalation_cost, eprr_costs, grand_total = cost_breakdown(
                pred_val, eprr, sst_pct, owners_pct, cont_pct, esc_pct
            )

            result = {"Project Name": project_name}
            result.update({c: payload.get(c, np.nan) for c in feat_cols})
            result[y_name] = round(pred_val, 2)

            for k, v in eprr_costs.items():
                result[f"{k} Cost"] = v

            result["SST Cost"] = sst_cost
            result["Owner's Cost"] = owners_cost
            result["Cost Contingency"] = contingency_cost
            result["Escalation & Inflation"] = escalation_cost
            result["Grand Total"] = grand_total

            st.session_state.predictions.setdefault(ds_name_pred, []).append(result)
            st.session_state.widget_nonce += 1
            toast("Prediction added to Results.")

            cA, cB, cC, cD, cE = st.columns(5)
            with cA:
                st.metric("Predicted (Base)", f"{currency_pred} {pred_val:,.2f}")
            with cB:
                st.metric("Owner's", f"{currency_pred} {owners_cost:,.2f}")
            with cC:
                st.metric("Contingency", f"{currency_pred} {contingency_cost:,.2f}")
            with cD:
                st.metric("Escalation", f"{currency_pred} {escalation_cost:,.2f}")
            with cE:
                st.metric("Grand Total (incl. SST)", f"{currency_pred} {grand_total:,.2f}")

        except Exception as e:
            st.error(f"Prediction failed: {e}")

    st.markdown('<h4 style="margin:0;color:#000;">Batch (Excel)</h4>', unsafe_allow_html=True)
    xls = st.file_uploader("Upload Excel for batch prediction", type=["xlsx"], key=f"batch_xlsx__{st.session_state.widget_nonce}")
    if xls:
        file_id = f"{xls.name}_{xls.size}_{ds_name_pred}_{target_col_pred}"
        if file_id not in st.session_state.processed_excel_files:
            try:
                batch_df = pd.read_excel(xls)
                missing = [c for c in feat_cols if c not in batch_df.columns]
                if missing:
                    st.error(f"Missing required columns in Excel: {missing}")
                else:
                    preds = pipe.predict(batch_df[feat_cols])

                    for i, row in batch_df.iterrows():
                        name = row.get("Project Name", f"Project {i+1}")

                        owners_cost, sst_cost, contingency_cost, escalation_cost, eprr_costs, grand_total = cost_breakdown(
                            float(preds[i]), eprr, sst_pct, owners_pct, cont_pct, esc_pct
                        )

                        entry = {"Project Name": name}
                        entry.update(row[feat_cols].to_dict())
                        entry[y_name] = round(float(preds[i]), 2)

                        for k, v in eprr_costs.items():
                            entry[f"{k} Cost"] = v
                        entry["SST Cost"] = sst_cost
                        entry["Owner's Cost"] = owners_cost
                        entry["Cost Contingency"] = contingency_cost
                        entry["Escalation & Inflation"] = escalation_cost
                        entry["Grand Total"] = grand_total

                        st.session_state.predictions.setdefault(ds_name_pred, []).append(entry)

                    st.session_state.processed_excel_files.add(file_id)
                    st.session_state.widget_nonce += 1
                    toast("Batch prediction complete.")
                    st.rerun()
            except Exception as e:
                st.error(f"Batch prediction failed: {e}")
        else:
            st.info("This batch file was already processed (history prevents duplicates).")

    # ========================= RESULTS / EXPORT ==============================
    st.divider()
    st.markdown('<h3 style="margin-top:0;color:#000;">📄 Results</h3>', unsafe_allow_html=True)

    ds_name_res = st.selectbox("Dataset", list(st.session_state.datasets.keys()), key="ds_results")
    preds = st.session_state.predictions.get(ds_name_res, [])

    st.markdown(f'<h4 style="margin:0;color:#000;">Project Entries</h4><p>{len(preds)} saved</p>', unsafe_allow_html=True)
    if preds:
        if st.button("🗑️ Delete all entries", key="delete_all_entries_btn"):
            st.session_state.predictions[ds_name_res] = []
            st.session_state.processed_excel_files = set()
            st.session_state.widget_nonce += 1
            toast("All entries removed.", "🗑️")
            st.rerun()

    st.markdown('<h4 style="margin:0;color:#000;">Summary Table & Export</h4><p>Download</p>', unsafe_allow_html=True)

    if preds:
        df_preds = pd.DataFrame(preds)
        df_disp = df_preds.copy()
        num_cols = df_disp.select_dtypes(include=[np.number]).columns
        for col in num_cols:
            df_disp[col] = df_disp[col].apply(lambda x: format_with_commas(x))
        st.dataframe(df_disp, use_container_width=True, height=420)

        bio_xlsx = io.BytesIO()
        df_preds.to_excel(bio_xlsx, index=False, engine="openpyxl")
        bio_xlsx.seek(0)

        metrics = st.session_state._last_metrics
        metrics_json = json.dumps(metrics if metrics else {"info": "No metrics"}, indent=2, default=float)

        zip_bio = io.BytesIO()
        with zipfile.ZipFile(zip_bio, "w", zipfile.ZIP_DEFLATED) as zf:
            zf.writestr(f"{ds_name_res}_predictions.xlsx", bio_xlsx.getvalue())
            zf.writestr(f"{ds_name_res}_metrics.json", metrics_json)
        zip_bio.seek(0)

        st.download_button(
            "⬇️ Download All (ZIP)",
            data=zip_bio.getvalue(),
            file_name=f"{ds_name_res}_capex_all.zip",
            mime="application/zip",
            key="download_zip_btn",
        )
    else:
        st.info("No data to export yet.")


# =======================================================================================
# EXPORT HELPERS (Excel / PPT)
# =======================================================================================
def _format_ws_money(ws, start_row=2):
    ws.freeze_panes = "A2"
    for c in range(1, ws.max_column + 1):
        ws.column_dimensions[get_column_letter(c)].width = 18
    for r in range(start_row, ws.max_row + 1):
        for c in range(1, ws.max_column + 1):
            cell = ws.cell(r, c)
            if isinstance(cell.value, (int, float)) and c >= 3:
                cell.number_format = "#,##0.00"


def create_project_excel_report_capex(project_name, proj, currency=""):
    output = io.BytesIO()
    comps_df = project_components_df(proj)

    if comps_df.empty:
        with pd.ExcelWriter(output, engine="openpyxl") as writer:
            pd.DataFrame({"Info": [f"No components for project {project_name}"]}).to_excel(
                writer, sheet_name="Summary", index=False
            )
        output.seek(0)
        return output

    totals = project_totals(proj)

    summary_df = comps_df.copy()
    summary_df.loc[len(summary_df)] = {
        "Component": "TOTAL",
        "Dataset": "",
        "Model": "",
        "Base CAPEX": totals["capex_sum"],
        "Owner's Cost": totals["owners"],
        "Contingency": totals["cont"],
        "Escalation": totals["esc"],
        "SST": totals["sst"],
        "Grand Total": totals["grand_total"],
    }

    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        summary_df.to_excel(writer, sheet_name="Summary", index=False)
        ws = writer.sheets["Summary"]
        _format_ws_money(ws)

        max_row = ws.max_row
        max_col = ws.max_column

        for col_idx in range(4, max_col + 1):
            col_letter = get_column_letter(col_idx)
            ws.conditional_formatting.add(
                f"{col_letter}2:{col_letter}{max_row-1}",
                ColorScaleRule(
                    start_type="percentile",
                    start_value=10,
                    start_color="FFE0F7FA",
                    mid_type="percentile",
                    mid_value=50,
                    mid_color="FF80DEEA",
                    end_type="percentile",
                    end_value=90,
                    end_color="FF00838F",
                ),
            )

        chart = BarChart()
        chart.title = "Grand Total by Component"
        data = Reference(ws, min_col=10, max_col=10, min_row=1, max_row=max_row - 1)
        cats = Reference(ws, min_col=1, min_row=2, max_row=max_row - 1)
        chart.add_data(data, titles_from_data=True)
        chart.set_categories(cats)
        chart.y_axis.title = f"Cost ({currency})".strip()
        chart.x_axis.title = "Component"
        chart.height = 10
        chart.width = 18
        ws.add_chart(chart, "L2")

        line = LineChart()
        line.title = "Base CAPEX Trend"
        data_capex = Reference(ws, min_col=4, max_col=4, min_row=1, max_row=max_row - 1)
        line.add_data(data_capex, titles_from_data=True)
        line.set_categories(cats)
        line.y_axis.title = f"Base CAPEX ({currency})".strip()
        line.height = 10
        line.width = 18
        ws.add_chart(line, "L20")

        comps_df.to_excel(writer, sheet_name="Components Detail", index=False)
        _format_ws_money(writer.sheets["Components Detail"])

    output.seek(0)
    return output


def create_project_pptx_report_capex(project_name, proj, currency=""):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    layout_title_only = prs.slide_layouts[5]
    layout_title_content = prs.slide_layouts[1]

    slide = prs.slides.add_slide(layout_title_only)
    title = slide.shapes.title
    title.text = f"CAPEX Project Report\n{project_name}"
    p = title.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 161, 155)

    comps_df = project_components_df(proj)
    comps = proj.get("components", [])
    totals = project_totals(proj)

    slide = prs.slides.add_slide(layout_title_content)
    slide.shapes.title.text = "Executive Summary"
    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    lines = [
        f"Project: {project_name}",
        f"Total Components: {len(comps)}",
        f"Total Base CAPEX: {currency} {totals['capex_sum']:,.2f}",
        f"Total Grand Total (incl. SST): {currency} {totals['grand_total']:,.2f}",
        "",
        "Components:",
    ]
    for c in comps:
        lines.append(f"• {c['component_type']}: {currency} {c['breakdown']['grand_total']:,.2f}")

    for i, line in enumerate(lines):
        para = body.paragraphs[0] if i == 0 else body.add_paragraph()
        para.text = line
        para.font.size = Pt(16)

    if not comps_df.empty:
        fig, ax = plt.subplots(figsize=(7, 4))
        ax.bar(comps_df["Component"], comps_df["Grand Total"])
        ax.set_title("Grand Total by Component")
        ax.set_ylabel(f"Cost ({currency})".strip())
        ax.tick_params(axis="x", rotation=25)
        ax.grid(axis="y", linestyle="--", alpha=0.4)
        fig.tight_layout()

        img_stream = io.BytesIO()
        fig.savefig(img_stream, format="png", dpi=200, bbox_inches="tight")
        plt.close(fig)
        img_stream.seek(0)

        slide = prs.slides.add_slide(layout_title_only)
        slide.shapes.title.text = "Grand Total by Component"
        slide.shapes.add_picture(img_stream, Inches(0.7), Inches(1.5), width=Inches(8.6))

        fig2, ax2 = plt.subplots(figsize=(7, 4))
        labels = comps_df["Component"]
        base = comps_df["Base CAPEX"]
        owners = comps_df["Owner's Cost"]
        cont = comps_df["Contingency"]
        esc = comps_df["Escalation"]
        sst = comps_df["SST"]

        bottom = np.zeros(len(labels))
        for vals, lab in [
            (base, "Base CAPEX"),
            (owners, "Owner"),
            (cont, "Contingency"),
            (esc, "Escalation"),
            (sst, "SST"),
        ]:
            ax2.bar(labels, vals, bottom=bottom, label=lab)
            bottom += np.array(vals, dtype=float)

        ax2.set_title("Cost Composition by Component")
        ax2.set_ylabel(f"Cost ({currency})".strip())
        ax2.tick_params(axis="x", rotation=25)
        ax2.grid(axis="y", linestyle="--", alpha=0.4)
        ax2.legend(fontsize=8, ncol=3)
        fig2.tight_layout()

        img_stream2 = io.BytesIO()
        fig2.savefig(img_stream2, format="png", dpi=200, bbox_inches="tight")
        plt.close(fig2)
        img_stream2.seek(0)

        slide2 = prs.slides.add_slide(layout_title_only)
        slide2.shapes.title.text = "Cost Composition by Component"
        slide2.shapes.add_picture(img_stream2, Inches(0.7), Inches(1.5), width=Inches(8.6))

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


def create_comparison_excel_report_capex(projects_dict, currency=""):
    output = io.BytesIO()

    summary_rows = []
    for name, proj in projects_dict.items():
        t = project_totals(proj)
        summary_rows.append(
            {
                "Project": name,
                "Components": len(proj.get("components", [])),
                "CAPEX Sum": t["capex_sum"],
                "Owner": t["owners"],
                "Contingency": t["cont"],
                "Escalation": t["esc"],
                "SST": t["sst"],
                "Grand Total": t["grand_total"],
            }
        )

    summary_df = pd.DataFrame(summary_rows)

    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        summary_df.to_excel(writer, sheet_name="Projects Summary", index=False)
        ws = writer.sheets["Projects Summary"]
        _format_ws_money(ws)

        max_row = ws.max_row
        max_col = ws.max_column

        for col_idx in range(3, max_col + 1):
            col_letter = get_column_letter(col_idx)
            ws.conditional_formatting.add(
                f"{col_letter}2:{col_letter}{max_row}",
                ColorScaleRule(
                    start_type="percentile",
                    start_value=10,
                    start_color="FFE3F2FD",
                    mid_type="percentile",
                    mid_value=50,
                    mid_color="FF90CAF9",
                    end_type="percentile",
                    end_value=90,
                    end_color="FF1565C0",
                ),
            )

        chart = BarChart()
        chart.title = "Grand Total by Project"
        data = Reference(ws, min_col=8, max_col=8, min_row=1, max_row=max_row)
        cats = Reference(ws, min_col=1, min_row=2, max_row=max_row)
        chart.add_data(data, titles_from_data=True)
        chart.set_categories(cats)
        chart.y_axis.title = f"Cost ({currency})".strip()
        chart.x_axis.title = "Project"
        chart.height = 10
        chart.width = 18
        ws.add_chart(chart, "J2")

        for name, proj in projects_dict.items():
            dfc = project_components_df(proj)
            if dfc.empty:
                continue
            sheet_name = name[:31]
            dfc.to_excel(writer, sheet_name=sheet_name, index=False)
            _format_ws_money(writer.sheets[sheet_name])

    output.seek(0)
    return output


def create_comparison_pptx_report_capex(projects_dict, currency=""):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    layout_title_only = prs.slide_layouts[5]

    slide = prs.slides.add_slide(layout_title_only)
    title = slide.shapes.title
    title.text = "CAPEX Project Comparison"
    p = title.text_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 161, 155)

    rows = []
    for name, proj in projects_dict.items():
        t = project_totals(proj)
        rows.append(
            {
                "Project": name,
                "CAPEX Sum": t["capex_sum"],
                "Owner": t["owners"],
                "Contingency": t["cont"],
                "Escalation": t["esc"],
                "SST": t["sst"],
                "Grand Total": t["grand_total"],
            }
        )
    df_proj = pd.DataFrame(rows)

    if not df_proj.empty:
        fig, ax = plt.subplots(figsize=(7, 4))
        ax.bar(df_proj["Project"], df_proj["Grand Total"])
        ax.set_title("Grand Total by Project")
        ax.set_ylabel(f"Cost ({currency})".strip())
        ax.tick_params(axis="x", rotation=25)
        ax.grid(axis="y", linestyle="--", alpha=0.4)
        fig.tight_layout()

        img_stream = io.BytesIO()
        fig.savefig(img_stream, format="png", dpi=200, bbox_inches="tight")
        plt.close(fig)
        img_stream.seek(0)

        slide = prs.slides.add_slide(layout_title_only)
        slide.shapes.title.text = "Grand Total by Project"
        slide.shapes.add_picture(img_stream, Inches(0.7), Inches(1.5), width=Inches(8.6))

        fig2, ax2 = plt.subplots(figsize=(7, 4))
        labels = df_proj["Project"]
        base = df_proj["CAPEX Sum"]
        owners = df_proj["Owner"]
        cont = df_proj["Contingency"]
        esc = df_proj["Escalation"]
        sst = df_proj["SST"]

        bottom = np.zeros(len(labels))
        for vals, lab in [
            (base, "Base CAPEX"),
            (owners, "Owner"),
            (cont, "Contingency"),
            (esc, "Escalation"),
            (sst, "SST"),
        ]:
            ax2.bar(labels, vals, bottom=bottom, label=lab)
            bottom += np.array(vals, dtype=float)

        ax2.set_title("Cost Composition by Project")
        ax2.set_ylabel(f"Cost ({currency})".strip())
        ax2.tick_params(axis="x", rotation=25)
        ax2.grid(axis="y", linestyle="--", alpha=0.4)
        ax2.legend(fontsize=8, ncol=3)
        fig2.tight_layout()

        img_stream2 = io.BytesIO()
        fig2.savefig(img_stream2, format="png", dpi=200, bbox_inches="tight")
        plt.close(fig2)
        img_stream2.seek(0)

        slide2 = prs.slides.add_slide(layout_title_only)
        slide2.shapes.title.text = "Cost Composition by Project"
        slide2.shapes.add_picture(img_stream2, Inches(0.7), Inches(1.5), width=Inches(8.6))

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


# ======================================================================================
# MONTE CARLO EXPORT HELPERS (Excel + PPT) — TRUE FAN (probability axis)
# ======================================================================================
def create_monte_carlo_excel_report(
    project_name: str,
    df_proj_mc: pd.DataFrame,
    df_curve: pd.DataFrame,
    df_fan_true: pd.DataFrame,
    df_tornado: pd.DataFrame,
    df_comp_sims: pd.DataFrame,
) -> io.BytesIO:
    out = io.BytesIO()
    with pd.ExcelWriter(out, engine="openpyxl") as writer:
        df_proj_mc.to_excel(writer, sheet_name="ProjectSims", index=False)
        df_curve.to_excel(writer, sheet_name="Curve", index=False)
        df_fan_true.to_excel(writer, sheet_name="Fan_True", index=False)
        df_tornado.to_excel(writer, sheet_name="Tornado", index=False)
        df_comp_sims.to_excel(writer, sheet_name="ComponentSims", index=False)

        for name in ["ProjectSims", "Curve", "Fan_True", "Tornado", "ComponentSims"]:
            ws = writer.sheets[name]
            ws.freeze_panes = "A2"
            for c in range(1, ws.max_column + 1):
                ws.column_dimensions[get_column_letter(c)].width = 22
    out.seek(0)
    return out


def create_monte_carlo_pptx_report(
    project_name: str,
    currency: str,
    baseline_gt: float,
    budget: float,
    p50: float,
    p80: float,
    p90: float,
    exceed_prob_pct: float,
    df_curve: pd.DataFrame,
    df_fan_true: pd.DataFrame,
    df_tornado: pd.DataFrame,
) -> io.BytesIO:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    layout_title_only = prs.slide_layouts[5]
    layout_title_content = prs.slide_layouts[1]

    # Slide 1: Title
    s1 = prs.slides.add_slide(layout_title_only)
    s1.shapes.title.text = f"Monte Carlo Report\n{project_name}"

    # Slide 2: Summary
    s2 = prs.slides.add_slide(layout_title_content)
    s2.shapes.title.text = "Executive Summary"
    tf = s2.shapes.placeholders[1].text_frame
    tf.clear()
    lines = [
        f"Baseline GT: {currency} {baseline_gt:,.2f}",
        f"Budget: {currency} {budget:,.2f}",
        f"P50: {currency} {p50:,.2f}",
        f"P80: {currency} {p80:,.2f}",
        f"P90: {currency} {p90:,.2f}",
        f"P(> Budget): {exceed_prob_pct:.1f}%",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)

    # Slide 3: Curve - REMOVE EXCEEDANCE LINE
    fig, ax = plt.subplots(figsize=(8.5, 4.5))
    ax.plot(df_curve["x"], df_curve["cdf"], label="CDF (P ≤ X)")
    # Removed exceedance line
    ax.axvline(p50, linestyle=":", label="P50")
    ax.axvline(p80, linestyle=":", label="P80")
    ax.axvline(p90, linestyle=":", label="P90")
    ax.axvline(budget, linestyle="-.", label="Budget")
    ax.set_title("CDF Curve")
    ax.set_xlabel(f"Project Grand Total ({currency})")
    ax.set_ylabel("Probability")
    ax.grid(True, linestyle="--", alpha=0.3)
    ax.legend(fontsize=9, ncol=2)
    fig.tight_layout()

    img = io.BytesIO()
    fig.savefig(img, format="png", dpi=200, bbox_inches="tight")
    plt.close(fig)
    img.seek(0)

    s3 = prs.slides.add_slide(layout_title_only)
    s3.shapes.title.text = "CDF Curve"
    s3.shapes.add_picture(img, Inches(0.7), Inches(1.5), width=Inches(8.6))

    # Slide 4: TRUE Fan (probability axis)
    fig2, ax2 = plt.subplots(figsize=(8.5, 4.5))
    ax2.fill_between(df_fan_true["prob"], df_fan_true["p10"], df_fan_true["p90"], alpha=0.25, label="P10–P90")
    ax2.fill_between(df_fan_true["prob"], df_fan_true["p20"], df_fan_true["p80"], alpha=0.35, label="P20–P80")
    ax2.fill_between(df_fan_true["prob"], df_fan_true["p40"], df_fan_true["p60"], alpha=0.45, label="P40–P60")
    ax2.plot(df_fan_true["prob"], df_fan_true["p50"], linewidth=2.5, label="P50")
    ax2.axhline(budget, linestyle="-.", label="Budget")
    ax2.set_title("TRUE Fan Chart (Confidence → Cost)")
    ax2.set_xlabel("Confidence Level (%)")
    ax2.set_ylabel(f"Cost ({currency})")
    ax2.grid(True, linestyle="--", alpha=0.3)
    ax2.legend(fontsize=9, ncol=2)
    fig2.tight_layout()

    img2 = io.BytesIO()
    fig2.savefig(img2, format="png", dpi=200, bbox_inches="tight")
    plt.close(fig2)
    img2.seek(0)

    s4 = prs.slides.add_slide(layout_title_only)
    s4.shapes.title.text = "TRUE Fan Chart"
    s4.shapes.add_picture(img2, Inches(0.7), Inches(1.5), width=Inches(8.6))

    # Slide 5: Tornado
    df_t = df_tornado.sort_values("variance_share_pct", ascending=True)

    fig3, ax3 = plt.subplots(figsize=(8.5, 4.5))
    ax3.barh(df_t["Component"], df_t["variance_share_pct"])
    ax3.set_title("Sensitivity Tornado (Variance Share)")
    ax3.set_xlabel("Variance Contribution (%)")
    ax3.grid(True, axis="x", linestyle="--", alpha=0.3)
    fig3.tight_layout()

    img3 = io.BytesIO()
    fig3.savefig(img3, format="png", dpi=200, bbox_inches="tight")
    plt.close(fig3)
    img3.seek(0)

    s5 = prs.slides.add_slide(layout_title_only)
    s5.shapes.title.text = "Sensitivity Tornado"
    s5.shapes.add_picture(img3, Inches(0.7), Inches(1.5), width=Inches(8.6))

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out


# =======================================================================================
# PROJECT BUILDER TAB
# =======================================================================================
with tab_pb:
    st.markdown('<h4 style="margin:0;color:#000;">Project Builder</h4><p>Assemble multi-component CAPEX projects</p>', unsafe_allow_html=True)

    if not st.session_state.datasets:
        st.info("No dataset. Go to **Data** tab to upload or load.")
        st.stop()

    colA, colB = st.columns([2, 1])
    with colA:
        new_project_name = st.text_input("New Project Name", placeholder="e.g., CAPEX 2026", key="pb_new_project_name")
    with colB:
        if new_project_name and new_project_name not in st.session_state.projects:
            if st.button("Create Project", key="pb_create_project_btn"):
                st.session_state.projects[new_project_name] = {"components": [], "totals": {}, "currency": ""}
                toast(f"Project '{new_project_name}' created.")
                st.rerun()

    if not st.session_state.projects:
        st.info("Create a project above, then add components.")
        st.stop()

    existing_projects = list(st.session_state.projects.keys())
    proj_sel = st.selectbox("Select project to work on", existing_projects, key="pb_project_select")

    ds_names = sorted(st.session_state.datasets.keys())
    dataset_for_comp = st.selectbox("Dataset for this component", ds_names, key="pb_dataset_for_component")
    df_comp = st.session_state.datasets[dataset_for_comp]

    target_col_comp = get_last_column_target(df_comp)
    st.session_state[f"target_col__{dataset_for_comp}"] = target_col_comp
    curr_ds = get_currency_symbol(df_comp, target_col_comp)

    default_label = st.session_state.component_labels.get(dataset_for_comp, "")
    component_type = st.text_input(
        "Component type (Asset / Scope)",
        value=(default_label or "Platform / Pipeline / Subsea / Well"),
        key=f"pb_component_type_{proj_sel}",
    )

    try:
        pipe_c, _, feat_cols_c, y_name_c, best_name_c = train_best_model_cached(
            df_comp,
            target_col_comp,
            test_size=0.2,
            random_state=42,
            dataset_key=dataset_for_comp,
        )
    except Exception as e:
        st.error(f"Component model setup failed: {e}")
        st.stop()

    st.markdown("**Component Feature Inputs (1 row)**")
    comp_input_key = f"pb_input_row__{proj_sel}__{dataset_for_comp}"
    if comp_input_key not in st.session_state:
        st.session_state[comp_input_key] = {c: np.nan for c in feat_cols_c}

    comp_row_df = pd.DataFrame([st.session_state[comp_input_key]], columns=feat_cols_c)
    comp_editor_key = f"pb_editor__{proj_sel}__{dataset_for_comp}__{st.session_state.widget_nonce}"
    comp_edited = st.data_editor(comp_row_df, num_rows="fixed", use_container_width=True, key=comp_editor_key)
    comp_payload = comp_edited.iloc[0].to_dict()

    st.markdown("---")
    st.markdown("**WBS Level 1**")
    cp1, cp2 = st.columns(2)
    with cp1:
        st.markdown("use +/-")
        eng_pb = st.number_input("Engineering", 0.0, 100.0, 0.0, 1.0, key=f"pb_eng_{proj_sel}")  # Changed to 0.0
        procurement_pb = st.number_input("Procurement", 0.0, 100.0, 0.0, 1.0, key=f"pb_procurement_{proj_sel}")  # Changed to 0.0
        fabrication_pb = st.number_input("Fabrication/Construction", 0.0, 100.0, 0.0, 1.0, key=f"pb_fabrication_{proj_sel}")  # Changed to 0.0
        ti_pb = st.number_input("Transportation & Installation (T&I)", 0.0, 100.0, 0.0, 1.0, key=f"pb_ti_{proj_sel}")  # Changed to 0.0

        eprr_pb = {"Engineering": eng_pb, "Procurement": procurement_pb, "Fabrication/Construction": fabrication_pb, "Transportation & Installation": ti_pb}
        eprr_total_pb = sum(eprr_pb.values())
        st.caption(f"WBS total: **{eprr_total_pb:.2f}%**")

    with cp2:
        st.markdown("use +/-")
        sst_pb = st.number_input("SST", 0.0, 100.0, 0.0, 0.5, key=f"pb_sst_{proj_sel}")
        owners_pb = st.number_input("Owner's Cost", 0.0, 100.0, 0.0, 0.5, key=f"pb_owners_{proj_sel}")
        cont_pb = st.number_input("Contingency", 0.0, 100.0, 0.0, 0.5, key=f"pb_cont_{proj_sel}")
        esc_pb = st.number_input("Escalation & Inflation", 0.0, 100.0, 0.0, 0.5, key=f"pb_esc_{proj_sel}")

    if st.button("➕ Predict & Add Component", key=f"pb_add_comp_{proj_sel}_{dataset_for_comp}"):
        try:
            base_pred = single_prediction(pipe_c, feat_cols_c, comp_payload)
            owners_cost, sst_cost, contingency_cost, escalation_cost, eprr_costs, grand_total = cost_breakdown(
                base_pred, eprr_pb, sst_pb, owners_pb, cont_pb, esc_pb
            )

            comp_entry = {
                "component_type": component_type or default_label or "Component",
                "dataset": dataset_for_comp,
                "model_used": best_name_c,
                "inputs": {k: comp_payload.get(k, np.nan) for k in feat_cols_c},
                "feature_cols": list(feat_cols_c),  # needed for MC
                "prediction": base_pred,
                "breakdown": {
                    "eprr_costs": eprr_costs,
                    "eprr_pct": eprr_pb,
                    "sst_cost": sst_cost,
                    "owners_cost": owners_cost,
                    "contingency_cost": contingency_cost,
                    "escalation_cost": escalation_cost,
                    "grand_total": grand_total,
                    "target_col": y_name_c,  # last col
                    "sst_pct": float(sst_pb),
                    "owners_pct": float(owners_pb),
                    "cont_pct": float(cont_pb),
                    "esc_pct": float(esc_pb),
                },
            }

            st.session_state.projects[proj_sel]["components"].append(comp_entry)
            st.session_state.component_labels[dataset_for_comp] = component_type or default_label
            st.session_state.projects[proj_sel]["currency"] = curr_ds

            st.session_state.widget_nonce += 1
            toast(f"Component added to project '{proj_sel}'.")
            st.rerun()
        except Exception as e:
            st.error(f"Failed to add component: {e}")

    st.markdown("---")
    st.markdown("### Current Project Overview")

    proj = st.session_state.projects[proj_sel]
    comps = proj.get("components", [])
    if not comps:
        st.info("No components yet. Add at least one above.")
        st.stop()

    dfc = project_components_df(proj)
    curr = proj.get("currency", "") or curr_ds

    st.dataframe(
        dfc.style.format(
            {
                "Base CAPEX": "{:,.2f}",
                "Owner's Cost": "{:,.2f}",
                "Contingency": "{:,.2f}",
                "Escalation": "{:,.2f}",
                "SST": "{:,.2f}",
                "Grand Total": "{:,.2f}",
            }
        ),
        use_container_width=True,
    )

    t = project_totals(proj)
    proj["totals"] = {"capex_sum": t["capex_sum"], "grand_total": t["grand_total"]}

    col_t1, col_t2, col_t3 = st.columns(3)
    with col_t1:
        st.metric("Project CAPEX (Base)", f"{curr} {t['capex_sum']:,.2f}")
    with col_t2:
        st.metric("Project SST", f"{curr} {t['sst']:,.2f}")
    with col_t3:
        st.metric("Project Grand Total (incl. SST)", f"{curr} {t['grand_total']:,.2f}")

    st.markdown("#### Component Cost Composition")
    df_cost = dfc[["Component", "Base CAPEX", "Owner's Cost", "Contingency", "Escalation", "SST"]].copy()
    df_cost = df_cost.rename(columns={"Base CAPEX": "CAPEX", "Owner's Cost": "Owner"})
    df_melt = df_cost.melt(id_vars="Component", var_name="Cost Type", value_name="Value")
    fig_stack = px.bar(
        df_melt,
        x="Component",
        y="Value",
        color="Cost Type",
        barmode="stack",
        labels={"Value": f"Cost ({curr})"},
    )
    st.plotly_chart(fig_stack, use_container_width=True)

    st.markdown("#### Components")
    for idx, c in enumerate(comps):
        col1, col2, col3 = st.columns([4, 2, 1])
        with col1:
            st.write(f"**{c['component_type']}** — *{c['dataset']}* — {c.get('model_used', 'N/A')}")
        with col2:
            st.write(f"Grand Total: {curr} {c['breakdown']['grand_total']:,.2f}")
        with col3:
            if st.button("🗑️", key=f"pb_del_comp_{proj_sel}_{idx}"):
                comps.pop(idx)
                st.session_state.widget_nonce += 1
                toast("Component removed.", "🗑️")
                st.rerun()

    st.markdown("---")
    st.markdown("#### Export / Import Project")

    col_dl1, col_dl2, col_dl3 = st.columns(3)

    with col_dl1:
        excel_report = create_project_excel_report_capex(proj_sel, proj, curr)
        st.download_button(
            "⬇️ Download Project Excel",
            data=excel_report,
            file_name=f"{proj_sel}_CAPEX_Report.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            key=f"pb_dl_excel_{proj_sel}",
        )

    with col_dl2:
        pptx_report = create_project_pptx_report_capex(proj_sel, proj, curr)
        st.download_button(
            "⬇️ Download Project PowerPoint",
            data=pptx_report,
            file_name=f"{proj_sel}_CAPEX_Report.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            key=f"pb_dl_ppt_{proj_sel}",
        )

    with col_dl3:
        st.download_button(
            "⬇️ Download Project (JSON)",
            data=json.dumps(proj, indent=2, default=float),
            file_name=f"{proj_sel}.json",
            mime="application/json",
            key=f"pb_dl_json_{proj_sel}",
        )

    up_json = st.file_uploader("Import project JSON", type=["json"], key=f"pb_import_{proj_sel}__{st.session_state.widget_nonce}")
    if up_json is not None:
        try:
            data = json.load(up_json)
            st.session_state.projects[proj_sel] = data
            st.session_state.widget_nonce += 1
            toast("Project imported successfully.")
            st.rerun()
        except Exception as e:
            st.error(f"Failed to import project JSON: {e}")


# =======================================================================================
# 🎲 MONTE CARLO TAB (Curve + TRUE Fan + Tornado + Exports)
# =======================================================================================
with tab_mc:
    st.markdown('<h3 style="margin-top:0;color:#000;">🎲 Monte Carlo</h3>', unsafe_allow_html=True)
    st.caption("Simulate uncertainty per component and roll-up to project Grand Total distribution.")

    if not st.session_state.projects:
        st.info("No projects found. Create a project first in **🏗️ Project Builder**.")
        st.stop()

    proj_names = list(st.session_state.projects.keys())
    proj_sel_mc = st.selectbox("Select project", proj_names, key="mc_project_select")

    proj = st.session_state.projects[proj_sel_mc]
    comps = proj.get("components", [])
    if not comps:
        st.warning("This project has no components. Add components in **🏗️ Project Builder** first.")
        st.stop()

    if not st.session_state.datasets:
        st.warning("No datasets loaded. Monte Carlo needs datasets to retrain cached models.")
        st.stop()

    t_mc = project_totals(proj)
    curr = proj.get("currency", "") or ""
    baseline_gt = float(t_mc["grand_total"])
    st.info(f"Baseline (current Project Grand Total): **{curr} {baseline_gt:,.2f}**")

    st.markdown("### Simulation settings")
    mcA, mcB, mcC = st.columns(3)
    with mcA:
        mc_n_sims = st.number_input("Simulations", 500, 50000, 5000, 500, key=f"mc_n_{proj_sel_mc}")
        mc_seed = st.number_input("Random seed", 0, 999999, 42, 1, key=f"mc_seed_{proj_sel_mc}")
    with mcB:
        mc_feat_sigma = st.slider("Feature uncertainty (±% stdev)", 0.0, 30.0, 5.0, 0.5, key=f"mc_feat_{proj_sel_mc}")
        mc_pct_sigma = st.slider("Percent uncertainty (± abs stdev)", 0.0, 10.0, 1.0, 0.1, key=f"mc_pct_{proj_sel_mc}")
    with mcC:
        mc_norm_eprr = st.checkbox("Normalize WBS to 100% each simulation", False, key=f"mc_norm_{proj_sel_mc}")
        mc_budget = st.number_input(
            "Budget threshold (Project Grand Total)",
            min_value=0.0,
            value=float(baseline_gt),
            step=1000.0,
            key=f"mc_budget_{proj_sel_mc}",
        )

    st.markdown("### Scenario buckets (Project Grand Total vs baseline)")
    sb1, sb2, sb3 = st.columns(3)
    with sb1:
        mc_low = st.slider("Low < baseline by (%)", 0, 50, 10, 1, key=f"mc_low_{proj_sel_mc}")
    with sb2:
        mc_band = st.slider("Base band ± (%)", 1, 50, 10, 1, key=f"mc_band_{proj_sel_mc}")
    with sb3:
        mc_high = st.slider("High > baseline by (%)", 0, 50, 10, 1, key=f"mc_high_{proj_sel_mc}")

    if st.button("Run Monte Carlo", type="primary", key=f"mc_run_{proj_sel_mc}"):
        try:
            with st.spinner("Running Monte Carlo for each component and rolling up..."):
                n = int(mc_n_sims)
                project_gt = np.zeros(n, dtype=float)

                comp_series = {}  # component_name -> np.array sims

                comp_summ_rows = []
                for idx, comp in enumerate(comps):
                    ds_name = comp["dataset"]
                    df_ds = st.session_state.datasets.get(ds_name)
                    if df_ds is None:
                        raise ValueError(f"Dataset not found in session: {ds_name}")

                    # target stored per component as "target_col" (which is last col)
                    target_col = comp["breakdown"].get("target_col")
                    if not target_col:
                        # fallback to last column now
                        target_col = get_last_column_target(df_ds)

                    pipe_tmp, _, feat_cols_tmp, target_used, _ = train_best_model_cached(
                        df_ds, target_col, test_size=0.2, random_state=42, dataset_key=ds_name
                    )

                    feat_cols = comp.get("feature_cols") or feat_cols_tmp
                    payload = comp.get("inputs") or {}

                    eprr_pct = comp["breakdown"].get("eprr_pct", {})
                    sst_pct_c = float(comp["breakdown"].get("sst_pct", 0.0))
                    owners_pct_c = float(comp["breakdown"].get("owners_pct", 0.0))
                    cont_pct_c = float(comp["breakdown"].get("cont_pct", 0.0))
                    esc_pct_c = float(comp["breakdown"].get("esc_pct", 0.0))

                    comp_seed = int(mc_seed) + (idx + 1) * 101

                    df_mc_c = monte_carlo_component(
                        model_pipe=pipe_tmp,
                        feature_cols=list(feat_cols),
                        base_payload=payload,
                        n_sims=n,
                        seed=comp_seed,
                        feature_sigma_pct=float(mc_feat_sigma),
                        pct_sigma_abs=float(mc_pct_sigma),
                        eprr=eprr_pct,
                        sst_pct=sst_pct_c,
                        owners_pct=owners_pct_c,
                        cont_pct=cont_pct_c,
                        esc_pct=esc_pct_c,
                        normalize_eprr_each_draw=bool(mc_norm_eprr),
                    )

                    c_name = comp["component_type"]
                    comp_series[c_name] = df_mc_c["grand_total"].to_numpy(dtype=float)
                    project_gt += comp_series[c_name]

                    comp_summ_rows.append(
                        {
                            "Component": c_name,
                            "Dataset": ds_name,
                            "P50": float(df_mc_c["grand_total"].quantile(0.50)),
                            "P80": float(df_mc_c["grand_total"].quantile(0.80)),
                            "P90": float(df_mc_c["grand_total"].quantile(0.90)),
                        }
                    )

                df_proj_mc = pd.DataFrame({"project_grand_total": project_gt})

                buckets, pct_delta = scenario_bucket_from_baseline(
                    df_proj_mc["project_grand_total"], baseline_gt, mc_low, mc_band, mc_high
                )
                df_proj_mc["Scenario"] = buckets
                df_proj_mc["%Δ vs baseline"] = pct_delta

                p50 = float(df_proj_mc["project_grand_total"].quantile(0.50))
                p80 = float(df_proj_mc["project_grand_total"].quantile(0.80))
                p90 = float(df_proj_mc["project_grand_total"].quantile(0.90))
                exceed_prob = float((df_proj_mc["project_grand_total"] > float(mc_budget)).mean()) * 100.0

                df_comp_mc = pd.DataFrame(comp_summ_rows)

                # -------------------------
                # Curve data (CDF only - no exceedance)
                # -------------------------
                df_curve = df_proj_mc[["project_grand_total"]].sort_values("project_grand_total").reset_index(drop=True)
                n_curve = len(df_curve)
                df_curve["cdf"] = np.arange(1, n_curve + 1) / n_curve
                df_curve["exceed"] = 1.0 - df_curve["cdf"]
                df_curve = df_curve.rename(columns={"project_grand_total": "x"})  # consistent

                # -------------------------
                # TRUE Fan data (probability axis)
                # -------------------------
                p_grid = np.linspace(0.01, 0.99, 200)
                vals_proj = df_proj_mc["project_grand_total"]
                df_fan_true = pd.DataFrame({
                    "prob": p_grid * 100.0,
                    "p10": np.full_like(p_grid, float(vals_proj.quantile(0.10))),
                    "p20": np.full_like(p_grid, float(vals_proj.quantile(0.20))),
                    "p40": np.full_like(p_grid, float(vals_proj.quantile(0.40))),
                    "p50": np.full_like(p_grid, float(vals_proj.quantile(0.50))),
                    "p60": np.full_like(p_grid, float(vals_proj.quantile(0.60))),
                    "p80": np.full_like(p_grid, float(vals_proj.quantile(0.80))),
                    "p90": np.full_like(p_grid, float(vals_proj.quantile(0.90))),
                })

                # -------------------------
                # Tornado (variance share per component)
                # -------------------------
                df_comp_sims = pd.DataFrame(comp_series)  # columns: components
                var_total = float(np.var(df_proj_mc["project_grand_total"].to_numpy(dtype=float), ddof=1))
                tornado_rows = []
                for c in df_comp_sims.columns:
                    v = float(np.var(df_comp_sims[c].to_numpy(dtype=float), ddof=1))
                    share = (v / var_total * 100.0) if var_total > 0 else 0.0
                    tornado_rows.append(
                        {
                            "Component": c,
                            "std_dev": float(np.std(df_comp_sims[c].to_numpy(dtype=float), ddof=1)),
                            "variance": v,
                            "variance_share_pct": share,
                            "mean": float(np.mean(df_comp_sims[c].to_numpy(dtype=float))),
                            "p50": float(np.quantile(df_comp_sims[c].to_numpy(dtype=float), 0.50)),
                            "p80": float(np.quantile(df_comp_sims[c].to_numpy(dtype=float), 0.80)),
                            "p90": float(np.quantile(df_comp_sims[c].to_numpy(dtype=float), 0.90)),
                        }
                    )
                df_tornado = pd.DataFrame(tornado_rows).sort_values("variance_share_pct", ascending=False)

            # -------------------------
            # UI OUTPUTS
            # -------------------------
            st.markdown("### Summary")
            m1, m2, m3, m4 = st.columns(4)
            m1.metric("P50 Project Grand Total", f"{curr} {p50:,.2f}")
            m2.metric("P80 Project Grand Total", f"{curr} {p80:,.2f}")
            m3.metric("P90 Project Grand Total", f"{curr} {p90:,.2f}")
            m4.metric("P(> Budget)", f"{exceed_prob:.1f}%")

            # Histogram
            fig_hist = px.histogram(df_proj_mc, x="project_grand_total", nbins=60, title="Project Grand Total distribution")
            st.plotly_chart(fig_hist, use_container_width=True)

            # CDF curve only (no exceedance) + budget marker
            fig_curve = go.Figure()
            fig_curve.add_trace(go.Scatter(x=df_curve["x"], y=df_curve["cdf"], mode="lines", name="CDF (P ≤ X)", line=dict(width=3)))
            # Removed exceedance line

            for val, label in [(p50, "P50"), (p80, "P80"), (p90, "P90")]:
                fig_curve.add_vline(x=val, line_width=2, line_dash="dot", annotation_text=label, annotation_position="top")

            fig_curve.add_vline(
                x=float(mc_budget),
                line_width=3,
                line_dash="dashdot",
                annotation_text=f"Budget | P(>Budget)={exceed_prob:.1f}%",
                annotation_position="top right",
            )

            fig_curve.update_layout(
                title="Monte Carlo Cost Curve (CDF)",
                xaxis_title=f"Project Grand Total ({curr})",
                yaxis_title="Probability",
                yaxis=dict(tickformat=".0%"),
                hovermode="x unified",
                height=480,
            )
            st.plotly_chart(fig_curve, use_container_width=True)

            # TRUE Fan chart (Project) — probability axis
            fig_fan = go.Figure()

            # P10–P90 band
            fig_fan.add_trace(go.Scatter(x=df_fan_true["prob"], y=df_fan_true["p90"], mode="lines", line=dict(width=0), showlegend=False))
            fig_fan.add_trace(go.Scatter(
                x=df_fan_true["prob"], y=df_fan_true["p10"],
                mode="lines", fill="tonexty",
                fillcolor="rgba(0,161,155,0.20)",
                name="P10–P90"
            ))

            # P20–P80 band
            fig_fan.add_trace(go.Scatter(x=df_fan_true["prob"], y=df_fan_true["p80"], mode="lines", line=dict(width=0), showlegend=False))
            fig_fan.add_trace(go.Scatter(
                x=df_fan_true["prob"], y=df_fan_true["p20"],
                mode="lines", fill="tonexty",
                fillcolor="rgba(108,77,211,0.25)",
                name="P20–P80"
            ))

            # P40–P60 band
            fig_fan.add_trace(go.Scatter(x=df_fan_true["prob"], y=df_fan_true["p60"], mode="lines", line=dict(width=0), showlegend=False))
            fig_fan.add_trace(go.Scatter(
                x=df_fan_true["prob"], y=df_fan_true["p40"],
                mode="lines", fill="tonexty",
                fillcolor="rgba(0,0,0,0.20)",
                name="P40–P60"
            ))

            # Median
            fig_fan.add_trace(go.Scatter(x=df_fan_true["prob"], y=df_fan_true["p50"], mode="lines", line=dict(width=3), name="P50"))
            fig_fan.add_hline(y=float(mc_budget), line_width=2, line_dash="dashdot", annotation_text="Budget", annotation_position="top left")

            fig_fan.update_layout(
                title="TRUE Fan Chart — Project Grand Total (Confidence → Cost)",
                xaxis_title="Confidence Level (%)",
                yaxis_title=f"Cost ({curr})",
                height=480,
            )
            st.plotly_chart(fig_fan, use_container_width=True)

            # Scenario buckets
            bucket_counts = df_proj_mc["Scenario"].value_counts().reset_index()
            bucket_counts.columns = ["Scenario", "Count"]
            fig_bucket = px.bar(bucket_counts, x="Scenario", y="Count", title="Scenario bucket counts")
            st.plotly_chart(fig_bucket, use_container_width=True)

            # Component summary
            st.markdown("### Component summary (P50/P80/P90)")
            st.dataframe(
                df_comp_mc.style.format({"P50": "{:,.2f}", "P80": "{:,.2f}", "P90": "{:,.2f}"}),
                use_container_width=True,
            )

            # TRUE Fan chart (Component)
            st.markdown("### TRUE Fan Chart (Component)")
            comp_pick = st.selectbox("Select component", list(df_comp_sims.columns), key=f"mc_comp_pick_{proj_sel_mc}")

            vals_c = pd.Series(df_comp_sims[comp_pick].to_numpy(dtype=float))
            df_fan_c = pd.DataFrame({
                "prob": p_grid * 100.0,
                "p10": np.full_like(p_grid, float(vals_c.quantile(0.10))),
                "p20": np.full_like(p_grid, float(vals_c.quantile(0.20))),
                "p40": np.full_like(p_grid, float(vals_c.quantile(0.40))),
                "p50": np.full_like(p_grid, float(vals_c.quantile(0.50))),
                "p60": np.full_like(p_grid, float(vals_c.quantile(0.60))),
                "p80": np.full_like(p_grid, float(vals_c.quantile(0.80))),
                "p90": np.full_like(p_grid, float(vals_c.quantile(0.90))),
            })

            fig_fan_c = go.Figure()
            fig_fan_c.add_trace(go.Scatter(x=df_fan_c["prob"], y=df_fan_c["p90"], mode="lines", line=dict(width=0), showlegend=False))
            fig_fan_c.add_trace(go.Scatter(x=df_fan_c["prob"], y=df_fan_c["p10"], fill="tonexty", mode="lines",
                                           fillcolor="rgba(0,161,155,0.20)", name="P10–P90"))
            fig_fan_c.add_trace(go.Scatter(x=df_fan_c["prob"], y=df_fan_c["p80"], mode="lines", line=dict(width=0), showlegend=False))
            fig_fan_c.add_trace(go.Scatter(x=df_fan_c["prob"], y=df_fan_c["p20"], fill="tonexty", mode="lines",
                                           fillcolor="rgba(108,77,211,0.25)", name="P20–P80"))
            fig_fan_c.add_trace(go.Scatter(x=df_fan_c["prob"], y=df_fan_c["p60"], mode="lines", line=dict(width=0), showlegend=False))
            fig_fan_c.add_trace(go.Scatter(x=df_fan_c["prob"], y=df_fan_c["p40"], fill="tonexty", mode="lines",
                                           fillcolor="rgba(0,0,0,0.20)", name="P40–P60"))
            fig_fan_c.add_trace(go.Scatter(x=df_fan_c["prob"], y=df_fan_c["p50"], mode="lines", line=dict(width=3), name="P50"))
            fig_fan_c.update_layout(
                title=f"TRUE Fan — Component: {comp_pick}",
                xaxis_title="Confidence Level (%)",
                yaxis_title=f"Cost ({curr})",
                height=460,
            )
            st.plotly_chart(fig_fan_c, use_container_width=True)

            # Tornado
            st.markdown("### Sensitivity Tornado (Variance Contribution)")
            fig_tornado = px.bar(
                df_tornado,
                x="variance_share_pct",
                y="Component",
                orientation="h",
                title="Sensitivity Tornado (Variance Share %)",
            )
            fig_tornado.update_layout(xaxis_title="Variance contribution (%)", height=520)
            st.plotly_chart(fig_tornado, use_container_width=True)

            st.dataframe(
                df_tornado.style.format(
                    {
                        "std_dev": "{:,.2f}",
                        "variance": "{:,.2f}",
                        "variance_share_pct": "{:,.1f}",
                        "mean": "{:,.2f}",
                        "p50": "{:,.2f}",
                        "p80": "{:,.2f}",
                        "p90": "{:,.2f}",
                    }
                ),
                use_container_width=True,
            )

            with st.expander("Show Monte Carlo table (first 200 rows)", expanded=False):
                st.dataframe(df_proj_mc.head(200), use_container_width=True)

            # Downloads (CSV + ZIP)
            st.markdown("### Download Monte Carlo results")
            csv_proj = df_proj_mc.to_csv(index=False).encode("utf-8")
            csv_comp = df_comp_mc.to_csv(index=False).encode("utf-8")

            d1, d2, d3 = st.columns(3)
            with d1:
                st.download_button(
                    "⬇️ Project MC (CSV)",
                    data=csv_proj,
                    file_name=f"{proj_sel_mc}_mc_project.csv",
                    mime="text/csv",
                    key=f"mc_dl_proj_{proj_sel_mc}",
                )
            with d2:
                st.download_button(
                    "⬇️ Component Summary (CSV)",
                    data=csv_comp,
                    file_name=f"{proj_sel_mc}_mc_components.csv",
                    mime="text/csv",
                    key=f"mc_dl_comp_{proj_sel_mc}",
                )

            bio_xlsx = io.BytesIO()
            with pd.ExcelWriter(bio_xlsx, engine="openpyxl") as writer:
                df_proj_mc.to_excel(writer, sheet_name="Project_MC", index=False)
                df_comp_mc.to_excel(writer, sheet_name="Component_Summary", index=False)
            bio_xlsx.seek(0)

            zip_bio = io.BytesIO()
            with zipfile.ZipFile(zip_bio, "w", zipfile.ZIP_DEFLATED) as zf:
                zf.writestr(f"{proj_sel_mc}_mc_project.csv", csv_proj)
                zf.writestr(f"{proj_sel_mc}_mc_components.csv", csv_comp)
                zf.writestr(f"{proj_sel_mc}_mc_results.xlsx", bio_xlsx.getvalue())
            zip_bio.seek(0)

            with d3:
                st.download_button(
                    "⬇️ MC Results (ZIP)",
                    data=zip_bio.getvalue(),
                    file_name=f"{proj_sel_mc}_mc_results.zip",
                    mime="application/zip",
                    key=f"mc_dl_zip_{proj_sel_mc}",
                )

            # Excel + PPT report exports (true fan + tornado + sims)
            st.markdown("### Download Monte Carlo reports (Excel / PPT)")
            mc_excel = create_monte_carlo_excel_report(
                project_name=proj_sel_mc,
                df_proj_mc=df_proj_mc.rename(columns={"project_grand_total": "Project Grand Total"}),
                df_curve=df_curve,
                df_fan_true=df_fan_true,
                df_tornado=df_tornado,
                df_comp_sims=df_comp_sims,
            )
            mc_ppt = create_monte_carlo_pptx_report(
                project_name=proj_sel_mc,
                currency=curr,
                baseline_gt=float(baseline_gt),
                budget=float(mc_budget),
                p50=float(p50),
                p80=float(p80),
                p90=float(p90),
                exceed_prob_pct=float(exceed_prob),
                df_curve=df_curve,
                df_fan_true=df_fan_true,
                df_tornado=df_tornado,
            )

            r1, r2 = st.columns(2)
            with r1:
                st.download_button(
                    "⬇️ Download MC Excel (Curve+Fan+Tornado+Sims)",
                    data=mc_excel,
                    file_name=f"{proj_sel_mc}_MonteCarlo_Report.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    key=f"mc_dl_excel_{proj_sel_mc}",
                )
            with r2:
                st.download_button(
                    "⬇️ Download MC PowerPoint (Curve+Fan+Tornado)",
                    data=mc_ppt,
                    file_name=f"{proj_sel_mc}_MonteCarlo_Report.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key=f"mc_dl_ppt_{proj_sel_mc}",
                )

        except Exception as e:
            st.error(f"Monte Carlo failed: {e}")


# =======================================================================================
# COMPARE PROJECTS TAB
# =======================================================================================
with tab_compare:
    st.markdown('<h4 style="margin:0;color:#000;">Compare Projects</h4><p>Portfolio-level CAPEX view</p>', unsafe_allow_html=True)

    proj_names = list(st.session_state.projects.keys())
    if len(proj_names) < 2:
        st.info("Create at least two projects in the Project Builder tab to compare.")
        st.stop()

    compare_sel = st.multiselect(
        "Select projects to compare",
        proj_names,
        default=proj_names[:2],
        key="compare_projects_sel",
    )

    if len(compare_sel) < 2:
        st.warning("Select at least two projects for a meaningful comparison.")
        st.stop()

    rows = []
    for p in compare_sel:
        proj = st.session_state.projects[p]
        t = project_totals(proj)
        proj["totals"] = {"capex_sum": t["capex_sum"], "grand_total": t["grand_total"]}
        rows.append(
            {
                "Project": p,
                "Components": len(proj.get("components", [])),
                "CAPEX Sum": t["capex_sum"],
                "Owner": t["owners"],
                "Contingency": t["cont"],
                "Escalation": t["esc"],
                "SST": t["sst"],
                "Grand Total": t["grand_total"],
                "Currency": proj.get("currency", ""),
            }
        )

    df_proj = pd.DataFrame(rows)
    st.dataframe(
        df_proj[["Project", "Components", "CAPEX Sum", "SST", "Grand Total"]].style.format(
            {"CAPEX Sum": "{:,.2f}", "SST": "{:,.2f}", "Grand Total": "{:,.2f}"}
        ),
        use_container_width=True,
    )

    st.markdown("#### Grand Total by Project")
    fig_gt = px.bar(df_proj, x="Project", y="Grand Total", text="Grand Total", barmode="group")
    fig_gt.update_traces(texttemplate="%{text:,.0f}", textposition="outside")
    st.plotly_chart(fig_gt, use_container_width=True)

    st.markdown("#### Stacked Cost Composition by Project")
    df_melt = df_proj.melt(
        id_vars=["Project"],
        value_vars=["CAPEX Sum", "Owner", "Contingency", "Escalation", "SST"],
        var_name="Cost Type",
        value_name="Value",
    )
    fig_comp = px.bar(df_melt, x="Project", y="Value", color="Cost Type", barmode="stack")
    st.plotly_chart(fig_comp, use_container_width=True)

    st.markdown("#### Component-Level Details")
    for p in compare_sel:
        proj = st.session_state.projects[p]
        comps = proj.get("components", [])
        if not comps:
            continue
        with st.expander(f"Project: {p}"):
            rows_c = []
            for c in comps:
                eprr_costs = c["breakdown"].get("eprr_costs", {})
                eprr_str = ", ".join(f"{k}: {v:,.0f}" for k, v in eprr_costs.items() if float(v) != 0)
                rows_c.append(
                    {
                        "Component": c["component_type"],
                        "Dataset": c["dataset"],
                        "Model": c.get("model_used", ""),
                        "Base CAPEX": c["prediction"],
                        "Owner": c["breakdown"]["owners_cost"],
                        "Contingency": c["breakdown"]["contingency_cost"],
                        "Escalation": c["breakdown"]["escalation_cost"],
                        "SST": c["breakdown"]["sst_cost"],
                        "Grand Total": c["breakdown"]["grand_total"],
                        "EPRR Costs": eprr_str,
                    }
                )
            df_compd = pd.DataFrame(rows_c)
            st.dataframe(
                df_compd.style.format(
                    {
                        "Base CAPEX": "{:,.2f}",
                        "Owner": "{:,.2f}",
                        "Contingency": "{:,.2f}",
                        "Escalation": "{:,.2f}",
                        "SST": "{:,.2f}",
                        "Grand Total": "{:,.2f}",
                    }
                ),
                use_container_width=True,
            )

    st.markdown("---")
    st.markdown("#### Download Comparison Reports")

    col_c1, col_c2 = st.columns(2)
    projects_to_export = {name: st.session_state.projects[name] for name in compare_sel}
    currency_comp = st.session_state.projects[compare_sel[0]].get("currency", "")

    with col_c1:
        excel_comp = create_comparison_excel_report_capex(projects_to_export, currency_comp)
        st.download_button(
            "⬇️ Download Comparison Excel",
            data=excel_comp,
            file_name="CAPEX_Projects_Comparison.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            key="cmp_dl_excel",
        )

    with col_c2:
        pptx_comp = create_comparison_pptx_report_capex(projects_to_export, currency_comp)
        st.download_button(
            "⬇️ Download Comparison PowerPoint",
            data=pptx_comp,
            file_name="CAPEX_Projects_Comparison.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            key="cmp_dl_ppt",
        )
